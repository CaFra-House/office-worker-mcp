"""Test del núcleo: render_document, pdf_tools (form fill, ocr, convert, manipulate), temas. (v0.2.0)"""
import os, shutil, sys, pathlib
import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from office_worker.core.templates import render_pdf
from office_worker.core.word import create_word, edit_word, mail_merge
from office_worker.core.excel import create_excel, edit_excel, csv_excel_convert
from office_worker.core.slides import create_pptx
from office_worker.core.themes import load_theme, THEMES
from office_worker.core.pdf import read_pdf, extract_tables, list_form_fields, pdf_preview, pdf_extract_structured
from office_worker.core.security import safe_out, safe_url_fetcher
from office_worker.core.pdf_tools import (
    fill_pdf_form,
    ocr_pdf,
    convert_office_to_pdf,
    manipulate_pdf,
    pdf_redact,
)
from office_worker.core.pdf_to_excel import pdf_to_excel
from office_worker.core.office_reader import read_office
from office_worker.core.diff import document_diff
from office_worker.core.compliance import scrub_metadata, protect_office, verify_pdf_signature
from office_worker.skills import resolve_packaged_skill, install_skill, list_packaged_skills

TPL = """<h1>{{ titulo }}</h1><p class="muted">{{ subtitulo }} · {{ fecha }}</p>
{% if tabla is defined and tabla %}{{ tabla }}{% endif %}"""


def test_render_pdf_and_themes(tmp_path):
    out = str(tmp_path / "out.pdf")
    data = {
        "titulo": "Informe de Prueba",
        "subtitulo": "The Office Worker v0.2.0",
        "fecha": "2026-09-05",
        "headers": ["Rubro", "Valor"],
        "rows": [["Ingresos", "$1.200.000"], ["Gastos", "$800.000"]],
    }
    headers = data["headers"]; rows = data["rows"]
    thead = "".join(f"<th>{h}</th>" for h in headers)
    body = "\n".join("<tr>" + "".join(f"<td>{c}</td>" for c in r) + "</tr>" for r in rows)
    data["tabla"] = f"<table><thead><tr>{thead}</tr></thead><tbody>{body}</tbody></table>"

    # Test con tema nuevo corporate-blue
    path = render_pdf(TPL, out, data=data, theme="corporate-blue")
    assert os.path.exists(path), "PDF no creado"
    size = os.path.getsize(path)
    assert size > 500, f"PDF demasiado chico: {size}"
    assert open(path, "rb").read(5) == b"%PDF-", "No es PDF válido"


def test_render_pdf_with_logo(tmp_path):
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow no instalado")

    logo_path = str(tmp_path / "logo.png")
    img = Image.new("RGB", (150, 40), color=(15, 76, 129))
    img.save(logo_path)

    out = str(tmp_path / "out_logo.pdf")
    path = render_pdf("<h1>Con Logo</h1>", out, theme="claro", logo=logo_path)
    assert os.path.exists(path)
    assert os.path.getsize(path) > 500
    assert open(path, "rb").read(5) == b"%PDF-"


def test_themes_catalog():
    for name in ["aden", "claro", "oscuro", "minimal", "corporate-blue"]:
        t = load_theme(name)
        assert t["name"] == name or name in t["name"]
        assert "primary" in t and "bg" in t and "text" in t


def test_fill_pdf_form(tmp_path):
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import black, white
        from pypdf import PdfReader
    except ImportError:
        pytest.skip("reportlab o pypdf no disponibles")

    form_pdf = str(tmp_path / "form.pdf")
    filled_pdf = str(tmp_path / "filled.pdf")

    # Crear PDF con AcroForm interactivo usando reportlab
    c = canvas.Canvas(form_pdf)
    c.drawString(100, 700, "Nombre:")
    form = c.acroForm
    form.textfield(name="nombre", tooltip="Nombre completo", x=150, y=695, width=200, height=20, borderColor=black, fillColor=white)
    c.showPage()
    c.save()

    assert os.path.exists(form_pdf)

    # Rellenar formulario
    res_path = fill_pdf_form(form_pdf, {"nombre": "Ana Gomez"}, filled_pdf)
    assert os.path.exists(res_path)
    assert open(res_path, "rb").read(5) == b"%PDF-"

    # Verificar campo relleno
    r = PdfReader(res_path)
    fields = r.get_fields()
    assert fields is not None and "nombre" in fields
    assert fields["nombre"].get("/V") == "Ana Gomez"

    # Verificar excepción si fields_dict está vacío
    with pytest.raises(ValueError):
        fill_pdf_form(form_pdf, {}, str(tmp_path / "err.pdf"))


def test_convert_office_to_pdf(tmp_path):
    if not (shutil.which("soffice") or shutil.which("libreoffice")):
        pytest.skip("soffice/LibreOffice no disponible en el sistema")

    docx_path = str(tmp_path / "doc.docx")
    pdf_out = str(tmp_path / "doc.pdf")

    create_word(docx_path, title="Prueba Conversión", subtitle="LibreOffice", blocks=[{"type": "p", "text": "Contenido para convertir"}])
    assert os.path.exists(docx_path)
    # Magic bytes PK para docx (zip)
    assert open(docx_path, "rb").read(2) == b"PK"

    res_path = convert_office_to_pdf(docx_path, pdf_out)
    assert os.path.exists(res_path)
    assert os.path.getsize(res_path) > 0
    assert open(res_path, "rb").read(5) == b"%PDF-"


def test_manipulate_pdf(tmp_path):
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        pytest.skip("pypdf no disponible")

    pdf1 = str(tmp_path / "p1.pdf")
    pdf2 = str(tmp_path / "p2.pdf")
    merged = str(tmp_path / "merged.pdf")
    extracted = str(tmp_path / "extracted.pdf")
    rotated = str(tmp_path / "rotated.pdf")

    w1 = PdfWriter(); w1.add_blank_page(width=100, height=100); w1.write(pdf1)
    w2 = PdfWriter(); w2.add_blank_page(width=100, height=100); w2.write(pdf2)

    # 1. Merge
    res_m = manipulate_pdf(operation="merge", out=merged, files=[pdf1, pdf2])
    assert os.path.exists(res_m)
    assert len(PdfReader(res_m).pages) == 2
    assert open(res_m, "rb").read(5) == b"%PDF-"

    # 2. Extract
    res_e = manipulate_pdf(operation="extract", out=extracted, input_path=merged, pages="1")
    assert os.path.exists(res_e)
    assert len(PdfReader(res_e).pages) == 1
    assert open(res_e, "rb").read(5) == b"%PDF-"

    # 3. Rotate
    res_r = manipulate_pdf(operation="rotate", out=rotated, input_path=extracted, angle=90)
    assert os.path.exists(res_r)
    r_doc = PdfReader(res_r)
    assert r_doc.pages[0].rotation == 90
    assert open(res_r, "rb").read(5) == b"%PDF-"


def test_ocr_pdf(tmp_path):
    if not shutil.which("tesseract"):
        pytest.skip("Tesseract OCR no instalado en el sistema")
    try:
        from PIL import Image, ImageDraw
        import pytesseract
    except ImportError:
        pytest.skip("PIL o pytesseract no disponibles")

    img_path = str(tmp_path / "test_ocr.png")
    pdf_out = str(tmp_path / "searchable.pdf")

    img = Image.new("RGB", (300, 80), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((20, 30), "HELLO WORLD", fill=(0, 0, 0))
    img.save(img_path)

    text = ocr_pdf(img_path, lang="eng", out=pdf_out)
    clean_text = text.replace(" ", "").upper()
    assert "HELLO" in clean_text or "HELLOWORLD" in clean_text
    assert os.path.exists(pdf_out)
    assert os.path.getsize(pdf_out) > 0
    assert open(pdf_out, "rb").read(5) == b"%PDF-"


def test_safe_out_security(tmp_path):
    # 1. Path traversal a sistema protegido debe lanzar PermissionError
    with pytest.raises(PermissionError):
        safe_out("/etc/cron.d/malicious.pdf")

    with pytest.raises(PermissionError):
        safe_out("../../../../etc/shadow")

    # 2. Ruta vacía o directorio existente
    with pytest.raises(ValueError):
        safe_out("")
    with pytest.raises(ValueError):
        safe_out(str(tmp_path))

    # 3. Base dir permitido (sandboxing)
    allowed_dir = str(tmp_path / "sandbox")
    os.makedirs(allowed_dir, exist_ok=True)
    valid_path = os.path.join(allowed_dir, "report.pdf")
    assert safe_out(valid_path, base_dir=allowed_dir) == valid_path

    # Fuera del sandbox debe fallar
    with pytest.raises(PermissionError):
        safe_out(str(tmp_path / "outside.pdf"), base_dir=allowed_dir)


def test_safe_url_fetcher_security():
    # 1. SSRF remoto bloqueado
    with pytest.raises(PermissionError):
        safe_url_fetcher("http://169.254.169.254/latest/meta-data")
    with pytest.raises(PermissionError):
        safe_url_fetcher("https://example.com/malicious.css")

    # 2. LFI a directorio del sistema protegido bloqueado
    with pytest.raises(PermissionError):
        safe_url_fetcher("file:///etc/passwd")
    with pytest.raises(PermissionError):
        safe_url_fetcher("file:///root/.ssh/id_rsa")


def test_render_pdf_encrypted(tmp_path):
    from pypdf import PdfReader
    out = str(tmp_path / "encrypted.pdf")
    path = render_pdf(
        "<h1>Confidencial</h1><p>Datos protegidos</p>",
        out,
        password="clave_secreta_99",
        theme="corporate-blue"
    )
    assert os.path.exists(path)
    reader = PdfReader(path)
    assert reader.is_encrypted, "El PDF generado debe estar cifrado"
    assert reader.decrypt("clave_secreta_99") > 0, "Debe descifrarse correctamente con la contraseña"


def test_manipulate_pdf_encrypted(tmp_path):
    from pypdf import PdfReader, PdfWriter
    p1 = str(tmp_path / "p1.pdf")
    p2 = str(tmp_path / "p2.pdf")
    w = PdfWriter()
    w.add_blank_page(width=100, height=100)
    with open(p1, "wb") as f: w.write(f)
    with open(p2, "wb") as f: w.write(f)

    out_merged = str(tmp_path / "merged_enc.pdf")
    res = manipulate_pdf(operation="merge", out=out_merged, files=[p1, p2], password="pwd_merge_123")
    assert os.path.exists(res)
    reader = PdfReader(res)
    assert reader.is_encrypted
    assert reader.decrypt("pwd_merge_123") > 0


def test_create_word_docxtpl(tmp_path):
    from docx import Document
    tpl_path = str(tmp_path / "template.docx")
    out_docx = str(tmp_path / "rendered.docx")

    # Crear plantilla inicial con variable docxtpl
    doc = Document()
    doc.add_paragraph("Hola {{ cliente }}, tu saldo es {{ saldo }}.")
    doc.save(tpl_path)

    # Rellenar con create_word
    res = create_word(
        out_docx,
        template_docx=tpl_path,
        context={"cliente": "Banco CAFRA", "saldo": "$1.500.000"}
    )
    assert os.path.exists(res)
    assert open(res, "rb").read(2) == b"PK"

    # Verificar que el texto fue reemplazado
    doc_res = Document(res)
    full_text = " ".join(p.text for p in doc_res.paragraphs)
    assert "Banco CAFRA" in full_text
    assert "$1.500.000" in full_text


def test_read_pdf_consolidated(tmp_path):
    out = str(tmp_path / "sample_read.pdf")
    render_pdf(
        "<h1>Titulo</h1><table><tr><th>A</th><th>B</th></tr><tr><td>1</td><td>2</td></tr></table>",
        out
    )
    # Lectura consolidada con flags
    res = read_pdf(out, extract_tables=True, list_forms=True)
    assert res.get("n_pages") >= 1
    assert "pages" in res
    assert "tables_by_page" in res
    assert "is_form" in res
    assert res["is_form"] is False


def test_pdf_ocr_edge_case_no_text(tmp_path):
    from pypdf import PdfWriter
    blank_pdf = str(tmp_path / "blank.pdf")
    w = PdfWriter()
    w.add_blank_page(width=100, height=100)
    with open(blank_pdf, "wb") as f: w.write(f)

    # OCR en PDF vacío con max_pages
    text = ocr_pdf(blank_pdf, lang="eng", max_pages=1)
    assert isinstance(text, str)
    assert text == ""


def test_fill_pdf_form_not_a_form(tmp_path):
    from pypdf import PdfWriter
    blank_pdf = str(tmp_path / "blank.pdf")
    w = PdfWriter()
    w.add_blank_page(width=100, height=100)
    with open(blank_pdf, "wb") as f: w.write(f)

    with pytest.raises(ValueError, match="no contiene campos de formulario interactivos"):
        fill_pdf_form(blank_pdf, {"campo": "valor"}, str(tmp_path / "out_err.pdf"))


def test_packaged_templates_pack(tmp_path):
    from office_worker.core.templates_pack import (
        list_packaged_templates,
        resolve_template_path,
        get_template_schema,
    )
    tpls = list_packaged_templates()
    names = [t["name"] for t in tpls]
    assert "acta_meeting" in names
    assert "informe_ejecutivo" in names
    assert "factura_simple" in names
    assert "carta_formal" in names
    assert "checklist_auditoria" in names
    assert len(names) == 5

    # resolve_template_path
    p = resolve_template_path("acta_meeting")
    assert os.path.exists(p)
    assert p.endswith("acta_meeting.docx")

    schema = get_template_schema("acta_meeting")
    assert "variables" in schema
    assert "asistentes" in schema["variables"]

    # create_word con plantilla empaquetada por nombre
    out_acta = str(tmp_path / "acta_out.docx")
    res = create_word(
        out_acta,
        template_docx="acta_meeting",
        context={
            "titulo": "Acta de Reunión de Directorio",
            "fecha": "2026-09-05",
            "hora": "10:30",
            "lugar": "Sede ADEN",
            "asistentes": [{"nombre": "Julio Cardozo", "rol": "Director"}],
            "puntos": [{"orden": 1, "tema": "V0.4.0", "discusion": "Aprobación"}],
            "acuerdos": [{"acuerdo": "Despliegue", "responsable": "JC", "fecha_limite": "2026-09-10"}],
            "firmas": [{"nombre": "Julio Cardozo", "cargo": "Director"}],
        }
    )
    assert os.path.exists(res)
    assert open(res, "rb").read(2) == b"PK"

    # Verificar las otras 4 plantillas
    for t_name, ctx in [
        ("informe_ejecutivo", {"titulo": "Inf", "subtitulo": "Sub", "fecha": "2026", "autor": "A", "resumen": "R", "secciones": [{"titulo": "S1", "contenido": "C1"}], "indicadores": [{"kpi": "K", "meta": "M", "actual": "A", "estado": "OK"}], "conclusiones": "Fin"}),
        ("factura_simple", {"numero_factura": "001", "fecha": "2026", "fecha_vencimiento": "2026", "emisor": {"nombre": "E", "cuit_nif": "1", "direccion": "D", "contacto": "C"}, "receptor": {"nombre": "R", "cuit_nif": "2", "direccion": "D2"}, "items": [{"descripcion": "Item 1", "cantidad": 1, "precio_unitario": "10", "subtotal": "10"}], "subtotal": "10", "iva": "2.1", "total": "12.1", "condiciones_pago": "Contado"}),
        ("carta_formal", {"lugar_fecha": "BA, 2026", "destinatario": {"nombre": "D", "cargo": "C", "organizacion": "O", "direccion": "Dir"}, "asunto": "Asunto", "saludo": "Estimado", "cuerpo": "Cuerpo", "despedida": "Saludos", "firma": {"nombre": "F", "cargo": "CF", "organizacion": "OF"}}),
        ("checklist_auditoria", {"titulo": "Check", "auditor": "Aud", "fecha": "2026", "alcance": "All", "items": [{"item": "1", "criterio": "Crit", "estado": "OK", "evidencia": "Evid", "observaciones": "Obs"}], "conclusion": "Aprobado"}),
    ]:
        out_p = str(tmp_path / f"{t_name}_out.docx")
        r = create_word(out_p, template_docx=t_name, context=ctx)
        assert os.path.exists(r)
        assert open(r, "rb").read(2) == b"PK"


def test_render_pdf_watermark_and_footers(tmp_path):
    import fitz
    out = str(tmp_path / "watermark_footer.pdf")
    path = render_pdf(
        "<h1>Documento Confidencial</h1><p>Texto interno de prueba.</p>",
        out,
        watermark_text="CONFIDENCIAL",
        footer_left="ID Documento: 98765",
        footer_right="Grupo CaFra 2026",
        page_numbers=True,
    )
    assert os.path.exists(path)
    assert open(path, "rb").read(5) == b"%PDF-"
    doc = fitz.open(path)
    page_text = doc[0].get_text()
    doc.close()
    assert "CONFIDENCIAL" in page_text
    assert "ID Documento: 98765" in page_text
    assert "Grupo CaFra 2026" in page_text
    assert "Página 1 de 1" in page_text


def test_sign_pdf_stamp_and_digital(tmp_path):
    import fitz
    from PIL import Image, ImageDraw
    from office_worker.core.pdf_tools import sign_pdf

    # 1. Crear PDF base
    pdf_in = str(tmp_path / "to_sign.pdf")
    render_pdf("<h1>Contrato de Prueba</h1><p>Cláusula 1: Validez de firma.</p>", pdf_in)

    # 2. Crear sello PNG
    stamp_png = str(tmp_path / "stamp.png")
    img = Image.new("RGBA", (180, 60), color=(255, 255, 255, 0))
    d = ImageDraw.Draw(img)
    d.rectangle([(2, 2), (178, 58)], outline=(0, 51, 102, 255), width=2)
    d.text((15, 15), "SELLO OFICIAL CAFRA", fill=(0, 51, 102, 255))
    img.save(stamp_png)

    # 3. Estampar sello visual
    pdf_visual = str(tmp_path / "signed_visual.pdf")
    res = sign_pdf(pdf_in, pdf_visual, sello_img_path=stamp_png)
    assert os.path.exists(res)
    assert open(res, "rb").read(5) == b"%PDF-"
    doc = fitz.open(res)
    images = doc[-1].get_images()
    doc.close()
    assert len(images) > 0, "Debe contener al menos una imagen estampada (fitz.get_images > 0)"

    # 4. Firma digital con certificado auto-firmado
    try:
        import datetime
        from cryptography import x509
        from cryptography.x509.oid import NameOID
        from cryptography.hazmat.primitives import hashes, serialization
        from cryptography.hazmat.primitives.asymmetric import rsa

        key = rsa.generate_private_key(public_exponent=65537, key_size=2048)
        subject = issuer = x509.Name([
            x509.NameAttribute(NameOID.COUNTRY_NAME, "AR"),
            x509.NameAttribute(NameOID.ORGANIZATION_NAME, "Grupo CaFra"),
            x509.NameAttribute(NameOID.COMMON_NAME, "Julio Cardozo"),
        ])
        cert = x509.CertificateBuilder().subject_name(
            subject
        ).issuer_name(
            issuer
        ).public_key(
            key.public_key()
        ).serial_number(
            x509.random_serial_number()
        ).not_valid_before(
            datetime.datetime.now(datetime.timezone.utc)
        ).not_valid_after(
            datetime.datetime.now(datetime.timezone.utc) + datetime.timedelta(days=1)
        ).sign(key, hashes.SHA256())

        cert_pem_path = str(tmp_path / "cert.pem")
        with open(cert_pem_path, "wb") as f:
            f.write(cert.public_bytes(serialization.Encoding.PEM))
            f.write(key.private_bytes(
                encoding=serialization.Encoding.PEM,
                format=serialization.PrivateFormat.PKCS8,
                encryption_algorithm=serialization.NoEncryption()
            ))

        pdf_digital = str(tmp_path / "signed_digital.pdf")
        res_d = sign_pdf(
            pdf_in,
            pdf_digital,
            sello_img_path=stamp_png,
            cert_pem=cert_pem_path,
            reason="Aprobación Técnica",
            location="Buenos Aires",
        )
        assert os.path.exists(res_d)
        assert open(res_d, "rb").read(5) == b"%PDF-"
        doc_d = fitz.open(res_d)
        assert len(doc_d[-1].get_images()) > 0
        doc_d.close()
    except ImportError:
        pass


def test_compress_pdf(tmp_path):
    import fitz
    from PIL import Image
    import io
    from office_worker.core.pdf_tools import compress_pdf

    # Crear PDF con imagen pesada
    img = Image.new("RGB", (1500, 1500), color=(180, 40, 40))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    png_bytes = buf.getvalue()

    large_pdf = str(tmp_path / "large.pdf")
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((50, 50), "PDF Grande para prueba de compresión")
    page.insert_image(fitz.Rect(50, 80, 450, 480), stream=png_bytes)
    doc.save(large_pdf)
    doc.close()

    size_orig = os.path.getsize(large_pdf)
    compressed_pdf = str(tmp_path / "compressed.pdf")

    res = compress_pdf(large_pdf, compressed_pdf, quality="med")
    assert res["status"] == "ok"
    assert os.path.exists(compressed_pdf)
    assert open(compressed_pdf, "rb").read(5) == b"%PDF-"
    assert res["size_after"] < size_orig
    assert res["savings_percent"] > 50

    # Verificar que conserva número de páginas y texto
    doc_c = fitz.open(compressed_pdf)
    assert len(doc_c) == 1
    assert "PDF Grande para prueba" in doc_c[0].get_text()
    doc_c.close()


def test_edit_excel(tmp_path):
    import openpyxl
    xlsx_in = str(tmp_path / "base.xlsx")
    create_excel(xlsx_in, title="Presupuesto", sheets=[{
        "name": "Finanzas",
        "headers": ["Item", "Monto"],
        "rows": [["Sueldos", 5000], ["Servidores", 1200]]
    }])
    assert os.path.exists(xlsx_in)

    # Nota: con title="Presupuesto", fila 1 es título, fila 2 encabezados, fila 3 y 4 datos.
    res = edit_excel(xlsx_in, operations=[
        {"op": "set_cell", "coordinate": "B3", "value": 5500},
        {"op": "append_row", "row": ["Licencias", 800, "=SUM(B3:B5)"]},
        {"op": "add_column", "header": "Moneda", "values": ["USD", "USD", "USD"]},
        {"op": "add_chart", "chart_type": "bar", "title": "Distribución", "target_cell": "E2"},
        {"op": "add_table", "table_style": "TableStyleMedium9"},
        {"op": "auto_filter"},
    ])
    assert res["status"] == "ok"
    assert res["fidelity"] == "rich"
    assert len(res["warnings"]) > 0  # aviso de fórmula
    assert os.path.exists(res["path"])
    assert open(res["path"], "rb").read(2) == b"PK"

    wb = openpyxl.load_workbook(res["path"])
    ws = wb["Finanzas"]
    assert ws["B3"].value == 5500
    assert ws.cell(5, 1).value == "Licencias"
    assert len(ws._charts) > 0
    wb.close()


def test_edit_word(tmp_path):
    from docx import Document
    docx_in = str(tmp_path / "base.docx")
    create_word(docx_in, title="Contrato Base", blocks=[
        {"type": "h1", "text": "Cláusula Primera"},
        {"type": "p", "text": "Este texto es el borrador viejo."},
    ])
    assert os.path.exists(docx_in)

    res = edit_word(docx_in, operations=[
        {"op": "append_paragraph", "text": "Párrafo final agregado.", "bold": True},
        {"op": "replace_text", "find": "viejo", "replace": "definitivo"},
        {"op": "insert_after_heading", "heading_text": "Cláusula Primera", "text": "Aclaración inmediata."},
        {"op": "append_table", "headers": ["Rol", "Nombre"], "rows": [["CEO", "Julio"]]},
    ])
    assert res["status"] == "ok"
    assert res["fidelity"] == "clean"
    assert len(res["warnings"]) > 0
    assert os.path.exists(res["path"])
    assert open(res["path"], "rb").read(2) == b"PK"

    doc = Document(res["path"])
    full_text = " ".join(p.text for p in doc.paragraphs)
    assert "definitivo" in full_text
    assert "viejo" not in full_text
    assert "Párrafo final agregado." in full_text
    assert "Aclaración inmediata." in full_text
    assert len(doc.tables) > 0


def test_read_pdf_extract_images(tmp_path):
    import fitz
    from PIL import Image
    import io

    img = Image.new("RGB", (100, 100), color=(0, 120, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    png_b = buf.getvalue()

    pdf_img = str(tmp_path / "doc_with_img.pdf")
    doc = fitz.open()
    p = doc.new_page(width=300, height=300)
    p.insert_image(fitz.Rect(20, 20, 120, 120), stream=png_b)
    doc.save(pdf_img)
    doc.close()

    out = read_pdf(pdf_img, extract_images=True, max_images=5)
    assert out["n_images"] == 1
    assert "images" in out
    assert out["images"][0]["data_url"].startswith("data:image/")
    assert out["images"][0]["width"] == 100


def test_pdf_to_excel(tmp_path):
    pdf_with_table = str(tmp_path / "table_doc.pdf")
    render_pdf("<h1>Tabla</h1><table><tr><th>Producto</th><th>Precio</th></tr><tr><td>Laptop</td><td>1500</td></tr></table>", pdf_with_table)

    xlsx_out = str(tmp_path / "from_pdf.xlsx")
    res = pdf_to_excel(pdf_with_table, xlsx_out)
    assert res["status"] == "ok"
    assert res["fidelity"] == "clean"
    assert res["n_tables"] >= 1
    assert os.path.exists(xlsx_out)
    assert open(xlsx_out, "rb").read(2) == b"PK"

    # Caso sin tablas -> lossy
    empty_pdf = str(tmp_path / "notable.pdf")
    render_pdf("<h1>Solo texto sin tablas</h1><p>Parrafo normal</p>", empty_pdf)
    xlsx_empty = str(tmp_path / "empty.xlsx")
    res_empty = pdf_to_excel(empty_pdf, xlsx_empty)
    assert res_empty["status"] == "ok"
    assert res_empty["fidelity"] == "lossy"
    assert len(res_empty["warnings"]) > 0


def test_read_office(tmp_path):
    from pptx import Presentation

    # 1. DOCX
    d_path = str(tmp_path / "doc_read.docx")
    create_word(d_path, title="Titulo DOCX", blocks=[{"type": "p", "text": "Parrafo DOCX"}])
    d_res = read_office(d_path)
    assert d_res["format"] == "docx"
    assert d_res["n_paragraphs"] >= 1
    assert "Titulo DOCX" in d_res["text"]

    # 2. XLSX (sin title para que headers estén en fila 1)
    x_path = str(tmp_path / "xls_read.xlsx")
    create_excel(x_path, sheets=[{"name": "H1", "headers": ["A", "B"], "rows": [[1, 2]]}])
    x_res = read_office(x_path)
    assert x_res["format"] == "xlsx"
    assert x_res["n_sheets"] == 1
    assert x_res["sheets"][0]["headers"] == ["A", "B"]

    # 3. PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Diapositiva Office Reader"
    p_path = str(tmp_path / "deck.pptx")
    prs.save(p_path)
    p_res = read_office(p_path)
    assert p_res["format"] == "pptx"
    assert p_res["n_slides"] == 1
    assert "Diapositiva Office Reader" in p_res["text"]


def test_create_excel_table_autofilter_charts(tmp_path):
    import openpyxl
    xlsx_chart = str(tmp_path / "chart_test.xlsx")
    create_excel(
        xlsx_chart,
        title="Dashboard",
        sheets=[{
            "name": "KPIs",
            "headers": ["Mes", "Ventas", "Gastos"],
            "rows": [["Enero", 100, 80], ["Febrero", 150, 90]],
            "charts": [{"type": "line", "title": "Tendencia", "target_cell": "E2"}],
            "table_style": "TableStyleMedium9",
            "auto_filter": True,
        }]
    )
    assert os.path.exists(xlsx_chart)
    wb = openpyxl.load_workbook(xlsx_chart)
    ws = wb["KPIs"]
    assert len(ws._charts) == 1
    assert len(ws.tables) == 1 or ws.auto_filter.ref is not None
    wb.close()


def test_pdf_preview(tmp_path):
    import base64

    pdf_path = str(tmp_path / "sample_preview.pdf")
    render_pdf(
        "<h1>Documento de Prueba Preview</h1><p>Verificando renderizado PNG con PyMuPDF.</p>",
        pdf_path,
    )
    assert os.path.exists(pdf_path)

    # 1. Preview en memoria (data_url base64) con dpi configurable (default 110)
    res = pdf_preview(pdf_path, dpi=110)
    assert res["status"] == "ok"
    assert res["n_pages"] >= 1
    assert res["rendered_pages"] >= 1
    assert len(res["pages"]) >= 1
    assert res["data_url"].startswith("data:image/png;base64,")

    # Decodificar y verificar magic bytes válidos de PNG
    b64_str = res["data_url"].split(",", 1)[1]
    png_bytes = base64.b64decode(b64_str)
    assert png_bytes[:8] == b"\x89PNG\r\n\x1a\n", "Magic bytes de PNG inválidos"

    # 2. Preview guardando archivo en disco
    out_png = str(tmp_path / "saved_preview.png")
    res_disk = pdf_preview(pdf_path, output=out_png, max_pages=1, dpi=96)
    assert res_disk["status"] == "ok"
    assert os.path.exists(out_png)
    assert os.path.getsize(out_png) > 0
    with open(out_png, "rb") as f:
        assert f.read(8) == b"\x89PNG\r\n\x1a\n"


def test_pdf_redact_core(tmp_path):
    import fitz
    pdf_in = str(tmp_path / "sensitive.pdf")
    pdf_out = str(tmp_path / "redacted.pdf")

    # Crear PDF con texto confidencial y coordenadas específicas
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 100), "El código de seguridad es CLAVE_CONFIDENCIAL_998877 de acceso.")
    page.insert_text((50, 200), "REG_TOP_SECRET en coordenadas 200.")
    doc.save(pdf_in)
    doc.close()

    assert os.path.exists(pdf_in)

    # Redactar por texto buscado Y por coordenadas simultáneamente
    res = pdf_redact(
        input_path=pdf_in,
        output=pdf_out,
        search_text="CLAVE_CONFIDENCIAL_998877",
        regions=[{"page": 1, "x0": 45, "y0": 185, "x1": 300, "y1": 215}],
        fill_color="black",
    )
    assert res["status"] == "ok"
    assert os.path.exists(pdf_out)
    assert res["redactions_count"] >= 2
    assert open(pdf_out, "rb").read(5) == b"%PDF-"

    # Verificación irreversible: el texto fue purgado del stream de fitz
    doc_out = fitz.open(pdf_out)
    extracted_text = doc_out[0].get_text()
    assert "CLAVE_CONFIDENCIAL_998877" not in extracted_text, "Texto confidencial aún presente tras redacción"
    assert "REG_TOP_SECRET" not in extracted_text, "Región de coordenadas aún presente tras redacción"
    doc_out.close()


def test_pdf_flatten_core(tmp_path):
    try:
        from reportlab.pdfgen import canvas
        from pypdf import PdfReader
        import fitz
    except ImportError:
        pytest.skip("reportlab, pypdf o fitz no disponible")

    form_pdf = str(tmp_path / "form_to_flatten.pdf")
    flat_pdf = str(tmp_path / "flattened.pdf")

    c = canvas.Canvas(form_pdf)
    c.drawString(50, 700, "Usuario Registrado:")
    c.acroForm.textfield(name="usuario", value="Juan Perez", x=180, y=695, width=150, height=20)
    c.showPage()
    c.save()

    r_before = PdfReader(form_pdf)
    assert r_before.get_fields() is not None, "El PDF de prueba debe tener campos AcroForm"

    # Aplanar formulario mediante manipulate_pdf con operation='flatten'
    res_path = manipulate_pdf(operation="flatten", input_path=form_pdf, out=flat_pdf)
    assert os.path.exists(res_path)
    assert open(res_path, "rb").read(5) == b"%PDF-"

    # Verificar que los campos interactivos ya no existen
    r_after = PdfReader(flat_pdf)
    assert not r_after.get_fields(), "Los campos de formulario no fueron aplanados (get_fields no está vacío)"

    # Verificar que el contenido visual sigue legible en el texto estático
    doc_flat = fitz.open(flat_pdf)
    text_content = doc_flat[0].get_text()
    assert "Juan Perez" in text_content or "Usuario Registrado" in text_content
    doc_flat.close()


def test_pdf_extract_structured_core(tmp_path):
    # Generar un PDF con tabla declarativa vía render_pdf
    pdf_in = str(tmp_path / "doc_with_table.pdf")
    html = """<h1>Reporte Financiero</h1>
    <table border="1">
        <tr><th>Concepto</th><th>Monto</th></tr>
        <tr><td>Suscripciones</td><td>50000</td></tr>
        <tr><td>Servicios</td><td>20000</td></tr>
    </table>"""
    render_pdf(html, pdf_in)

    assert os.path.exists(pdf_in)

    # 1. Extracción a JSON
    json_out = str(tmp_path / "extracted.json")
    res_json = pdf_extract_structured(pdf_in, format="json", output=json_out)
    assert res_json["status"] == "ok"
    assert res_json["format"] == "json"
    assert res_json["n_pages"] >= 1
    assert len(res_json["pages"]) >= 1
    assert res_json["n_tables"] >= 1
    tbl = res_json["pages"][0]["tables"][0]
    assert len(tbl) >= 2
    assert "Concepto" in tbl[0] or "Monto" in tbl[0]
    assert os.path.exists(json_out)

    # 2. Extracción a Markdown legible con tablas pipe
    md_out = str(tmp_path / "extracted.md")
    res_md = pdf_extract_structured(pdf_in, format="markdown", output=md_out)
    assert res_md["status"] == "ok"
    assert res_md["format"] == "markdown"
    assert "## Página 1" in res_md["content"]
    assert "| Concepto | Monto |" in res_md["content"] or "Concepto" in res_md["content"]
    assert "| --- | --- |" in res_md["content"]
    assert os.path.exists(md_out)


def test_pdf_split_smart_core(tmp_path):
    import fitz
    concat_pdf = str(tmp_path / "concatenated.pdf")
    split_target = str(tmp_path / "split_doc.pdf")

    # Crear PDF concatenado: Doc 1 (pág 0) + Página blanca separadora (pág 1) + Doc 2 (pág 2)
    doc = fitz.open()
    p0 = doc.new_page()
    p0.insert_text((50, 100), "DOCUMENTO 1: Contrato de Prestacion de Servicios")
    doc.new_page()  # Página en blanco (sin texto ni imágenes)
    p2 = doc.new_page()
    p2.insert_text((50, 100), "DOCUMENTO 2: Acta de Directorio Separada")
    doc.save(concat_pdf)
    doc.close()

    assert os.path.exists(concat_pdf)

    # Ejecutar split_smart
    res = manipulate_pdf(operation="split_smart", input_path=concat_pdf, out=split_target)
    assert res["status"] == "ok"
    assert res["n_splits"] == 2
    assert len(res["files"]) == 2

    # Verificar que cada archivo es un PDF independiente válido
    for f_path in res["files"]:
        assert os.path.exists(f_path)
        assert os.path.getsize(f_path) > 0
        with open(f_path, "rb") as f:
            assert f.read(5) == b"%PDF-", f"Archivo {f_path} no es PDF válido"

    doc1 = fitz.open(res["files"][0])
    assert "DOCUMENTO 1" in doc1[0].get_text()
    assert "DOCUMENTO 2" not in doc1[0].get_text()
    doc1.close()

    doc2 = fitz.open(res["files"][1])
    assert "DOCUMENTO 2" in doc2[0].get_text()
    assert "DOCUMENTO 1" not in doc2[0].get_text()
    doc2.close()


def test_skills_packaging_and_cli(tmp_path):
    # 1. Resolver skill empaquetada oficial
    path_skill = resolve_packaged_skill("office-worker")
    assert path_skill.exists()
    assert "The Office Worker" in path_skill.read_text(encoding="utf-8")

    # 2. Catálogo de skills
    catalog = list_packaged_skills()
    assert len(catalog) >= 2
    names = [s["name"] for s in catalog]
    assert "office-worker" in names and "google-drive-gmail" in names

    # 3. Instalación idempotente en directorio temporal
    dest = tmp_path / "hermes_skills"
    res = install_skill("office-worker", dest_dir=dest)
    assert res["status"] == "ok"
    installed_file = dest / "office-worker" / "SKILL.md"
    assert installed_file.exists()
    assert installed_file.stat().st_size > 500

    # 4. CLI con HOME fake sin tocar entornos reales
    from pathlib import Path
    from office_worker.cli import main as cli_main
    fake_home = str(tmp_path / "fake_user_home")
    old_home = os.environ.get("HOME")
    try:
        os.environ["HOME"] = fake_home
        ret = cli_main(["skill", "install", "office-worker"])
        assert ret == 0
        cli_installed = Path(fake_home) / ".hermes" / "skills" / "office-worker" / "SKILL.md"
        assert cli_installed.exists()
    finally:
        if old_home is not None:
            os.environ["HOME"] = old_home
        else:
            os.environ.pop("HOME", None)


def test_mail_merge_core(tmp_path):
    from docx import Document

    # 1. Crear plantilla .docx con placeholders {{ nombre }} y {{ ciudad }}
    tpl_file = str(tmp_path / "plantilla_bienvenida.docx")
    doc_tpl = Document()
    doc_tpl.add_heading("Carta de Bienvenida", level=1)
    doc_tpl.add_paragraph("Estimado/a {{ nombre }}, le damos la bienvenida a la sede de {{ ciudad }}.")
    doc_tpl.save(tpl_file)
    assert os.path.exists(tpl_file)

    # 2. Dataset CSV con 3 filas
    csv_file = str(tmp_path / "usuarios.csv")
    csv_content = "nombre,ciudad\nAlice,Madrid\nBob,Buenos Aires\nCarlos,Lima\n"
    with open(csv_file, "w", encoding="utf-8") as f:
        f.write(csv_content)

    # 3. Ejecutar mail_merge
    prefix = str(tmp_path / "carta_generada")
    res = mail_merge(template_path=tpl_file, dataset_csv=csv_file, output_prefix=prefix)

    assert res["status"] == "ok"
    assert res["n_docs"] == 3
    assert len(res["paths"]) == 3
    assert "nombre" in res["fields"] and "ciudad" in res["fields"]

    # 4. Verificar cada docx: PK magic bytes, tamaño y sustitución de valores única por fila
    expected_pairs = [
        ("Alice", "Madrid"),
        ("Bob", "Buenos Aires"),
        ("Carlos", "Lima"),
    ]

    for idx, (path, (expected_nombre, expected_ciudad)) in enumerate(zip(res["paths"], expected_pairs), start=1):
        assert os.path.exists(path)
        assert os.path.getsize(path) > 300
        with open(path, "rb") as f:
            assert f.read(2) == b"PK", f"Documento {path} no tiene magic bytes PK de DOCX/ZIP"

        doc = Document(path)
        full_text = " ".join(p.text for p in doc.paragraphs)
        assert expected_nombre in full_text, f"{expected_nombre} no encontrado en doc {idx}"
        assert expected_ciudad in full_text, f"{expected_ciudad} no encontrado en doc {idx}"

        # Verificar que no contiene valores de las otras filas
        for other_n, other_c in expected_pairs:
            if other_n != expected_nombre:
                assert other_n not in full_text, f"Valor de otra fila '{other_n}' presente en doc {idx}"

    # 5. Prueba adicional con dataset JSON y filtrado de campos
    json_dataset = [
        {"nombre": "Diana", "ciudad": "Bogotá", "extra": "omitir"},
        {"nombre": "Esteban", "ciudad": "Santiago", "extra": "omitir"},
    ]
    res_json = mail_merge(
        template_path=tpl_file,
        dataset_json=json_dataset,
        output_prefix=str(tmp_path / "json_carta"),
        fields=["nombre", "ciudad"],
    )
    assert res_json["status"] == "ok"
    assert res_json["n_docs"] == 2
    assert "extra" not in res_json["fields"]


def test_read_office_markdown_and_json(tmp_path):
    from docx import Document
    import openpyxl

    # 1. DOCX con Heading + Tabla -> formato markdown y backward compatibility
    docx_file = str(tmp_path / "informe_test.docx")
    doc = Document()
    doc.add_heading("Resumen Ejecutivo 2026", level=1)
    doc.add_paragraph("Este es un párrafo introductorio con detalles de gestión.")
    tbl = doc.add_table(rows=2, cols=2)
    tbl.rows[0].cells[0].text = "Indicador"
    tbl.rows[0].cells[1].text = "Resultado"
    tbl.rows[1].cells[0].text = "EBITDA"
    tbl.rows[1].cells[1].text = "+24%"
    doc.save(docx_file)

    # Formato Markdown
    res_md = read_office(docx_file, format="markdown")
    assert res_md["status"] == "ok"
    assert res_md["format"] == "markdown"
    assert "# Resumen Ejecutivo 2026" in res_md["content"]
    assert "| Indicador | Resultado |" in res_md["content"]
    assert "| --- | --- |" in res_md["content"]
    assert "| EBITDA | +24% |" in res_md["content"]

    # Formato JSON (default backward compatibility)
    res_json = read_office(docx_file)
    assert res_json["status"] == "ok"
    assert res_json["format"] == "docx"
    assert res_json["n_paragraphs"] >= 2
    assert res_json["n_tables"] == 1
    assert "Resumen Ejecutivo 2026" in res_json["text"]

    # 2. XLSX con tabla -> formato markdown por hoja
    xlsx_file = str(tmp_path / "metricas_test.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Ventas_Q1"
    ws.append(["Producto", "Ingresos"])
    ws.append(["Suscripcion", 50000])
    ws.append(["Servicios", 25000])
    wb.save(xlsx_file)

    res_xlsx_md = read_office(xlsx_file, format="markdown")
    assert res_xlsx_md["status"] == "ok"
    assert res_xlsx_md["format"] == "markdown"
    assert "## Ventas_Q1" in res_xlsx_md["content"]
    assert "| Producto | Ingresos |" in res_xlsx_md["content"]
    assert "| --- | --- |" in res_xlsx_md["content"]
    assert "| Suscripcion | 50000 |" in res_xlsx_md["content"]

    # XLSX formato default JSON
    res_xlsx_json = read_office(xlsx_file)
    assert res_xlsx_json["status"] == "ok"
    assert res_xlsx_json["format"] == "xlsx"
    assert res_xlsx_json["n_sheets"] == 1
    assert res_xlsx_json["sheets"][0]["name"] == "Ventas_Q1"


def test_create_pptx_with_charts_core(tmp_path):
    from pptx import Presentation

    pptx_file = str(tmp_path / "presentacion_graficos.pptx")
    slides = [
        {
            "title": "Portada Corporativa",
            "kicker": "The Office Worker v0.7.0",
            "bullets": ["Agenda", "Métricas del Ejercicio"],
        },
        {
            "title": "Ventas Trimestrales",
            "kicker": "Rendimiento",
            "chart": {
                "type": "bar",
                "title": "Facturación por Trimestre",
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "values": [120, 160, 210, 280],
            },
        },
        {
            "title": "Distribución de Mercado",
            "chart": {
                "type": "pie",
                "title": "Participación",
                "categories": ["Producto A", "Producto B", "Otros"],
                "values": [50, 30, 20],
            },
        },
    ]

    out_path = create_pptx(pptx_file, slides=slides, theme="corporate-blue")
    assert os.path.exists(out_path)
    assert os.path.getsize(out_path) > 3000

    # Recargar presentación y verificar que contiene gráficos DrawingML nativos
    prs = Presentation(out_path)
    assert len(prs.slides) == 3

    # Diapositiva 2 debe tener gráfico nativo (GraphicFrame con has_chart=True)
    slide2 = prs.slides[1]
    chart_shapes = [s for s in slide2.shapes if s.has_chart]
    assert len(chart_shapes) >= 1, "La diapositiva 2 no contiene un shape de tipo gráfico"
    chart_shape = chart_shapes[0]
    assert chart_shape.__class__.__name__ == "GraphicFrame", f"Shape no es GraphicFrame: {type(chart_shape)}"
    assert chart_shape.chart.has_title
    assert "Facturación por Trimestre" in chart_shape.chart.chart_title.text_frame.text

    # Diapositiva 3 debe tener gráfico de torta nativo
    slide3 = prs.slides[2]
    pie_shapes = [s for s in slide3.shapes if s.has_chart]
    assert len(pie_shapes) >= 1, "La diapositiva 3 no contiene un shape de tipo gráfico"
    assert pie_shapes[0].__class__.__name__ == "GraphicFrame"


def test_csv_excel_convert_roundtrip_core(tmp_path):
    # 1. Crear archivo CSV de prueba con variedad de tipos y posibles ambigüedades
    csv_in = str(tmp_path / "datos_origen.csv")
    csv_content = (
        "id,producto,precio,stock,codigo_postal\n"
        "1,Servidor Cloud,1500.5,3,01425\n"
        "2,Licencia MCP,250.0,10,01428\n"
        "3,Soporte Anual,800.0,5,01430\n"
    )
    with open(csv_in, "w", encoding="utf-8") as f:
        f.write(csv_content)

    xlsx_out = str(tmp_path / "convertido.xlsx")

    # 2. Conversión CSV -> XLSX
    res_xlsx = csv_excel_convert(csv_in, xlsx_out, direction="csv_to_xlsx", sheet="Inventario")
    assert res_xlsx["status"] == "ok"
    assert os.path.exists(res_xlsx["path"])
    assert res_xlsx["n_rows"] == 4
    assert res_xlsx["n_cols"] == 5

    # Verificar magic bytes PK del archivo Excel generado
    with open(res_xlsx["path"], "rb") as f:
        assert f.read(2) == b"PK", "El archivo Excel generado no tiene cabecera ZIP/PK"

    # Verificar presencia de advertencia honesta sobre código con ceros a la izquierda
    assert any("01425" in w for w in res_xlsx.get("warnings", [])), "Falta warning honesto sobre ceros a la izquierda"

    # 3. Conversión XLSX -> CSV (Roundtrip)
    csv_roundtrip = str(tmp_path / "roundtrip.csv")
    res_csv = csv_excel_convert(res_xlsx["path"], csv_roundtrip, direction="xlsx_to_csv")
    assert res_csv["status"] == "ok"
    assert os.path.exists(res_csv["path"])
    assert res_csv["n_rows"] == 4

    with open(csv_roundtrip, "r", encoding="utf-8") as f:
        rt_content = f.read()

    # Comprobar que no se perdieron datos numéricos ni strings con ceros
    assert "01425" in rt_content, "Se perdió el código con ceros a la izquierda '01425'"
    assert "1500.5" in rt_content or "1500.50" in rt_content, "Se corrompió el float 1500.5"
    assert "Servidor Cloud" in rt_content
    assert "3" in rt_content

    # 4. Multi-sheet export con sheet='all'
    import openpyxl
    wb = openpyxl.Workbook()
    ws1 = wb.active; ws1.title = "Ventas"; ws1.append(["A", "1"]); ws1.append(["B", "2"])
    ws2 = wb.create_sheet("Gastos"); ws2.append(["C", "3"]); ws2.append(["D", "4"])
    multi_xlsx = str(tmp_path / "multi.xlsx")
    wb.save(multi_xlsx)

    res_all = csv_excel_convert(multi_xlsx, str(tmp_path / "export_all.csv"), direction="xlsx_to_csv", sheet="all")
    assert res_all["status"] == "ok"
    assert res_all["n_files"] == 2
    for f_csv in res_all["files"]:
        assert os.path.exists(f_csv)


def test_document_diff_core(tmp_path):
    # 1. Crear documento A y B en formato .docx
    doc_a = str(tmp_path / "doc_a.docx")
    doc_b = str(tmp_path / "doc_b.docx")

    create_word(doc_a, title="Acuerdo Marco", blocks=[
        {"type": "p", "text": "Cláusula 1: Definiciones generales."},
        {"type": "p", "text": "Cláusula 2: Obligaciones a ser suprimidas."},
        {"type": "p", "text": "Cláusula 3: Jurisdicción y ley aplicable."},
    ])

    create_word(doc_b, title="Acuerdo Marco", blocks=[
        {"type": "p", "text": "Cláusula 1: Definiciones generales."},
        {"type": "p", "text": "Cláusula 2: Nueva obligación de auditoría agregada."},
        {"type": "p", "text": "Cláusula 3: Jurisdicción y ley aplicable."},
    ])

    # 2. Diff en JSON
    diff_json = document_diff(doc_a, doc_b, format="json")
    assert diff_json["status"] == "ok"
    assert diff_json["has_changes"] is True
    assert diff_json["summary"]["modified"] == 1
    assert diff_json["summary"]["unchanged"] >= 2
    assert any("textual" in w.lower() for w in diff_json["warnings"])
    assert any("legal-grade" in w.lower() for w in diff_json["warnings"])

    # 3. Diff con párrafo agregado y eliminado por separado
    doc_c = str(tmp_path / "doc_c.docx")
    create_word(doc_c, title="Acuerdo Marco", blocks=[
        {"type": "p", "text": "Cláusula 1: Definiciones generales."},
        {"type": "p", "text": "Cláusula 3: Jurisdicción y ley aplicable."},
        {"type": "p", "text": "Cláusula 4: Anexo confidencial nuevo."},
    ])
    diff_del_add = document_diff(doc_a, doc_c, format="markdown")
    assert diff_del_add["status"] == "ok"
    assert diff_del_add["summary"]["deleted"] == 1
    assert diff_del_add["summary"]["added"] == 1
    assert "diff_markdown" in diff_del_add
    assert "[Deleted]" in diff_del_add["diff_markdown"]
    assert "[Added]" in diff_del_add["diff_markdown"]

    # 4. Documentos idénticos
    diff_same = document_diff(doc_a, doc_a)
    assert diff_same["has_changes"] is False
    assert diff_same["summary"]["added"] == 0
    assert diff_same["summary"]["deleted"] == 0
    assert diff_same["summary"]["modified"] == 0


def test_scrub_metadata_core(tmp_path):
    import fitz
    import docx
    import openpyxl
    from pptx import Presentation

    # 1. Scrub PDF
    pdf_in = str(tmp_path / "doc_meta.pdf")
    render_pdf("<h1>Documento con Autor</h1>", pdf_in)
    doc_pdf = fitz.open(pdf_in)
    doc_pdf.set_metadata({"author": "Julio Cardozo", "title": "Confidencial", "subject": "Auditoria"})
    doc_pdf.save(pdf_in + ".tmp")
    doc_pdf.close()
    os.replace(pdf_in + ".tmp", pdf_in)

    pdf_out = str(tmp_path / "doc_scrubbed.pdf")
    res_pdf = scrub_metadata(pdf_in, pdf_out)
    assert res_pdf["status"] == "ok"
    assert os.path.exists(res_pdf["path"])
    assert open(res_pdf["path"], "rb").read(5) == b"%PDF-"
    doc_check = fitz.open(res_pdf["path"])
    assert doc_check.metadata.get("author") == ""
    assert doc_check.metadata.get("title") == ""
    doc_check.close()

    # 2. Scrub DOCX
    docx_in = str(tmp_path / "doc_meta.docx")
    d = docx.Document()
    d.add_paragraph("Texto del contrato con autoría.")
    d.core_properties.author = "Julio Cardozo"
    d.core_properties.title = "Contrato Secreto"
    d.core_properties.last_modified_by = "Elena Rostova"
    d.save(docx_in)

    docx_out = str(tmp_path / "doc_scrubbed.docx")
    res_docx = scrub_metadata(docx_in, docx_out)
    assert res_docx["status"] == "ok"
    assert os.path.exists(res_docx["path"])
    assert open(res_docx["path"], "rb").read(2) == b"PK"
    d_check = docx.Document(res_docx["path"])
    assert d_check.core_properties.author == ""
    assert d_check.core_properties.last_modified_by == ""
    assert d_check.core_properties.title == ""
    assert len(d_check.paragraphs) == 1
    assert "Texto del contrato con autoría." in d_check.paragraphs[0].text

    # 3. Scrub XLSX
    xlsx_in = str(tmp_path / "sheet_meta.xlsx")
    wb = openpyxl.Workbook()
    wb.properties.creator = "Julio Cardozo"
    wb.properties.title = "Balance"
    wb.save(xlsx_in)

    xlsx_out = str(tmp_path / "sheet_scrubbed.xlsx")
    res_xlsx = scrub_metadata(xlsx_in, xlsx_out)
    assert res_xlsx["status"] == "ok"
    wb_check = openpyxl.load_workbook(res_xlsx["path"])
    assert wb_check.properties.creator is None or wb_check.properties.creator == ""
    wb_check.close()

    # 4. Scrub PPTX
    pptx_in = str(tmp_path / "slide_meta.pptx")
    prs = Presentation()
    prs.core_properties.author = "Julio Cardozo"
    prs.save(pptx_in)

    pptx_out = str(tmp_path / "slide_scrubbed.pptx")
    res_pptx = scrub_metadata(pptx_in, pptx_out)
    assert res_pptx["status"] == "ok"
    prs_check = Presentation(res_pptx["path"])
    assert prs_check.core_properties.author == ""


def test_protect_office_core(tmp_path):
    import msoffcrypto
    import openpyxl
    import docx
    from zipfile import BadZipFile

    # 1. Proteger XLSX con contraseña
    xlsx_in = str(tmp_path / "financiero.xlsx")
    create_excel(xlsx_in, title="Finanzas Q3", sheets=[
        {"name": "Resumen", "headers": ["Concepto", "Total"], "rows": [["Ventas", 50000]]}
    ])

    xlsx_protected = str(tmp_path / "financiero_protegido.xlsx")
    res = protect_office(xlsx_in, xlsx_protected, password="MiClaveSegura2026!")
    assert res["status"] == "ok"
    assert os.path.exists(res["path"])
    assert res["encrypted"] is True
    assert res["bytes"] > 0

    # Abrir sin clave DEBE fallar
    with pytest.raises(BadZipFile):
        openpyxl.load_workbook(res["path"])

    # Descifrar con clave correcta DEBE tener éxito
    with open(res["path"], "rb") as f:
        of = msoffcrypto.OfficeFile(f)
        assert of.is_encrypted() is True
        of.load_key(password="MiClaveSegura2026!")
        dec_path = str(tmp_path / "financiero_descifrado.xlsx")
        with open(dec_path, "wb") as df:
            of.decrypt(df)

    wb_dec = openpyxl.load_workbook(dec_path)
    # Row 1 is Title, Row 2 is Headers, Row 3 is Data
    assert wb_dec["Resumen"]["A3"].value == "Ventas"
    assert wb_dec["Resumen"]["B3"].value == 50000
    wb_dec.close()

    # 2. Proteger DOCX
    docx_in = str(tmp_path / "contrato.docx")
    create_word(docx_in, title="Contrato", blocks=[{"type": "p", "text": "Cláusula secreta."}])
    docx_protected = str(tmp_path / "contrato_protegido.docx")
    res_docx = protect_office(docx_in, docx_protected, password="DocxPass2026!")
    assert res_docx["status"] == "ok"

    with pytest.raises(Exception):
        docx.Document(docx_protected)


def test_verify_pdf_signature_core(tmp_path):
    import datetime
    from cryptography import x509
    from cryptography.x509.oid import NameOID
    from cryptography.hazmat.primitives import hashes, serialization
    from cryptography.hazmat.primitives.asymmetric import rsa
    from office_worker.core.pdf_tools import sign_pdf

    # 1. PDF sin firma
    pdf_unsigned = str(tmp_path / "no_firma.pdf")
    render_pdf("<h1>Documento sin Firma</h1>", pdf_unsigned)
    res_un = verify_pdf_signature(pdf_unsigned)
    assert res_un["status"] == "ok"
    assert res_un["has_signature"] is False
    assert res_un["valid"] is False
    assert any("no digital signatures" in w.lower() for w in res_un["warnings"])

    # 2. Generar certificado y clave PEM
    key = rsa.generate_private_key(public_exponent=65537, key_size=2048)
    subject = issuer = x509.Name([
        x509.NameAttribute(NameOID.COMMON_NAME, "Julio Cardozo Auditor"),
    ])
    cert = (
        x509.CertificateBuilder()
        .subject_name(subject)
        .issuer_name(issuer)
        .public_key(key.public_key())
        .serial_number(x509.random_serial_number())
        .not_valid_before(datetime.datetime.now(datetime.timezone.utc))
        .not_valid_after(datetime.datetime.now(datetime.timezone.utc) + datetime.timedelta(days=1))
        .sign(key, hashes.SHA256())
    )

    cert_pem_path = str(tmp_path / "cert_audit.pem")
    with open(cert_pem_path, "wb") as f:
        f.write(cert.public_bytes(serialization.Encoding.PEM))
        f.write(key.private_bytes(
            encoding=serialization.Encoding.PEM,
            format=serialization.PrivateFormat.PKCS8,
            encryption_algorithm=serialization.NoEncryption()
        ))

    # 3. Firmar PDF con sign_pdf
    pdf_signed = str(tmp_path / "con_firma.pdf")
    sign_pdf(
        pdf_unsigned,
        pdf_signed,
        cert_pem=cert_pem_path,
        reason="Aprobación Auditoría v0.8.0",
        location="Buenos Aires",
    )
    assert os.path.exists(pdf_signed)

    # 4. Verificar firma con verify_pdf_signature
    res_ver = verify_pdf_signature(pdf_signed)
    assert res_ver["status"] == "ok"
    assert res_ver["has_signature"] is True
    assert res_ver["valid"] is True
    assert res_ver["intact"] is True
    assert res_ver["signatures_count"] >= 1
    assert "Julio Cardozo Auditor" in (res_ver["signer"] or "")
    assert res_ver["reason"] == "Aprobación Auditoría v0.8.0"
    assert res_ver["location"] == "Buenos Aires"
    assert any("self-signed" in w.lower() for w in res_ver["warnings"])


def test_sign_pdf_auto_generate_test_cert_and_verify(tmp_path):
    """Test obligatorio v0.8.1: sign_pdf con auto_generate_test_cert -> verify_pdf_signature."""
    from office_worker.core.pdf_tools import sign_pdf
    from PIL import Image

    pdf_raw = str(tmp_path / "base_raw.pdf")
    render_pdf("<h1>Documento para Firma Automática</h1>", pdf_raw)
    assert os.path.exists(pdf_raw)

    # 1. Sin firma: verify_pdf_signature DEBE retornar has_signature False y valid False
    r_unsigned = verify_pdf_signature(pdf_raw)
    assert r_unsigned["status"] == "ok"
    assert r_unsigned["has_signature"] is False
    assert r_unsigned["valid"] is False

    # 2. Sello visual solo (sin cert ni auto_generate_test_cert): NO debe tener firma criptográfica
    sello_png = str(tmp_path / "sello_solo.png")
    img = Image.new("RGBA", (120, 40), color=(0, 0, 0, 0))
    img.save(sello_png)

    pdf_visual_only = str(tmp_path / "visual_only.pdf")
    sign_pdf(pdf_raw, pdf_visual_only, sello_img_path=sello_png, auto_generate_test_cert=False)
    assert os.path.exists(pdf_visual_only)
    r_visual = verify_pdf_signature(pdf_visual_only)
    assert r_visual["has_signature"] is False
    assert r_visual["valid"] is False

    # 3. sign_pdf con auto_generate_test_cert=True
    pdf_signed_auto = str(tmp_path / "signed_auto_test.pdf")
    res_path = sign_pdf(
        pdf_raw,
        pdf_signed_auto,
        sello_img_path=sello_png,
        reason="Aprobación Automatizada de Test",
        location="CABA",
        auto_generate_test_cert=True,
    )
    assert os.path.exists(res_path)
    assert os.path.getsize(res_path) > 0
    assert open(res_path, "rb").read(5) == b"%PDF-"

    # 4. verify_pdf_signature sobre PDF firmado con auto_generate_test_cert
    r_signed = verify_pdf_signature(pdf_signed_auto)
    assert r_signed["status"] == "ok"
    assert r_signed["has_signature"] is True
    assert r_signed["valid"] is True
    assert r_signed["intact"] is True
    assert r_signed["signatures_count"] >= 1
    assert "Office Worker Test Certificate (Non-Production)" in (r_signed["signer"] or "")
    assert r_signed["reason"] == "Aprobación Automatizada de Test"
    assert r_signed["location"] == "CABA"
    assert any("self-signed" in w.lower() for w in r_signed["warnings"])


def test_next_steps_on_all_primary_tools(tmp_path):
    """Test obligatorio v0.8.1: verificar que las 8 tools principales incluyan next_steps con lista no vacía [max 2 EN]."""
    import json
    from office_worker.mcp_server import (
        create_word,
        create_excel,
        create_pptx,
        convert_to_pdf,
        render_document,
        sign_pdf,
        pdf_redact,
        mail_merge,
    )

    d = str(tmp_path)

    # 1. render_document
    r_rd = render_document("<h1>Test Next Steps</h1>", f"{d}/doc_render.pdf")
    assert r_rd["status"] == "ok"
    assert "next_steps" in r_rd and isinstance(r_rd["next_steps"], list) and 1 <= len(r_rd["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_rd["next_steps"])

    # 2. create_word
    r_cw = create_word(f"{d}/doc_word.docx", title="Test Word")
    assert r_cw["status"] == "ok"
    assert "next_steps" in r_cw and isinstance(r_cw["next_steps"], list) and 1 <= len(r_cw["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_cw["next_steps"])

    # 3. create_excel
    r_ce = create_excel(f"{d}/doc_excel.xlsx", title="Test Excel")
    assert r_ce["status"] == "ok"
    assert "next_steps" in r_ce and isinstance(r_ce["next_steps"], list) and 1 <= len(r_ce["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_ce["next_steps"])

    # 4. create_pptx
    r_cp = create_pptx(f"{d}/doc_slides.pptx", slides_json=json.dumps([{"title": "Slide 1"}]))
    assert r_cp["status"] == "ok"
    assert "next_steps" in r_cp and isinstance(r_cp["next_steps"], list) and 1 <= len(r_cp["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_cp["next_steps"])

    # 5. convert_to_pdf (si LibreOffice está disponible)
    if shutil.which("soffice") or shutil.which("libreoffice"):
        r_conv = convert_to_pdf(r_cw["path"], f"{d}/doc_converted.pdf")
        assert r_conv["status"] == "ok"
        assert "next_steps" in r_conv and isinstance(r_conv["next_steps"], list) and 1 <= len(r_conv["next_steps"]) <= 2
        assert all(isinstance(s, str) and len(s) > 0 for s in r_conv["next_steps"])

    # 6. sign_pdf (con auto_generate_test_cert=True)
    r_sign = sign_pdf(r_rd["path"], f"{d}/doc_signed.pdf", auto_generate_test_cert=True)
    assert r_sign["status"] == "ok"
    assert "next_steps" in r_sign and isinstance(r_sign["next_steps"], list) and 1 <= len(r_sign["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_sign["next_steps"])

    # 7. pdf_redact
    r_redact = pdf_redact(input_path=r_rd["path"], output=f"{d}/doc_redacted.pdf", search_text="Test")
    assert r_redact["status"] == "ok"
    assert "next_steps" in r_redact and isinstance(r_redact["next_steps"], list) and 1 <= len(r_redact["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_redact["next_steps"])

    # 8. mail_merge
    from docx import Document
    mm_tpl = f"{d}/tpl_mm.docx"
    doc_mm = Document(); doc_mm.add_paragraph("Hola {{ nombre }}"); doc_mm.save(mm_tpl)
    mm_csv = f"{d}/data_mm.csv"
    with open(mm_csv, "w", encoding="utf-8") as f:
        f.write("nombre\nAna\n")
    r_mm = mail_merge(template_path=mm_tpl, dataset_csv=mm_csv, output_prefix=f"{d}/mm_out")
    assert r_mm["status"] == "ok"
    assert "next_steps" in r_mm and isinstance(r_mm["next_steps"], list) and 1 <= len(r_mm["next_steps"]) <= 2
    assert all(isinstance(s, str) and len(s) > 0 for s in r_mm["next_steps"])


def test_excel_pivot_table_core(tmp_path):
    """Test obligatorio v0.9.0 (A): add_pivot en create_excel/edit_excel con agregados numéricos verificados."""
    import openpyxl
    xlsx_path = str(tmp_path / "ventas_corporativas.xlsx")

    # 1. Crear libro base con dataset de ventas
    create_excel(
        xlsx_path,
        sheets=[{
            "name": "Ventas",
            "headers": ["Vendedor", "Region", "Mes", "Monto"],
            "rows": [
                ["Ana", "Norte", "Q1", 1000],
                ["Beto", "Norte", "Q1", 1500],
                ["Carlos", "Sur", "Q1", 2000],
                ["Ana", "Norte", "Q2", 1200],
                ["Beto", "Sur", "Q2", 800],
                ["Carlos", "Sur", "Q2", 2200],
            ],
        }]
    )
    assert os.path.exists(xlsx_path)

    # 2. Agregar tabla dinámica con op add_pivot (Region x Monto SUM)
    res_edit = edit_excel(xlsx_path, operations=[
        {
            "op": "add_pivot",
            "sheet": "Ventas",
            "rows": "Region",
            "values": "Monto",
            "agg": "sum",
            "pivot_sheet": "Pivot_Region",
        },
        {
            "op": "add_pivot",
            "sheet": "Ventas",
            "rows": "Region",
            "cols": "Mes",
            "values": "Monto",
            "agg": "sum",
            "pivot_sheet": "Pivot_Region_Mes",
        },
    ])
    assert res_edit["status"] == "ok"
    assert "Pivot_Region" in res_edit["sheets_modified"]
    assert "Pivot_Region_Mes" in res_edit["sheets_modified"]

    # 3. Verificar en disco que las hojas pivot existen con formato y datos numéricos exactos
    wb = openpyxl.load_workbook(xlsx_path)
    assert "Pivot_Region" in wb.sheetnames
    assert "Pivot_Region_Mes" in wb.sheetnames

    # Hoja 1: Pivot_Region
    ws1 = wb["Pivot_Region"]
    # Fila 1 = Encabezados ('Region', 'Monto')
    assert ws1.cell(1, 1).value == "Region"
    assert ws1.cell(1, 2).value == "Monto"
    # Fila 2 = Norte: 1000 + 1500 + 1200 = 3700
    assert ws1.cell(2, 1).value == "Norte"
    assert ws1.cell(2, 2).value == 3700
    # Fila 3 = Sur: 2000 + 800 + 2200 = 5000
    assert ws1.cell(3, 1).value == "Sur"
    assert ws1.cell(3, 2).value == 5000
    # Autofilter activo
    assert ws1.auto_filter.ref is not None

    # Hoja 2: Pivot_Region_Mes (filas Region, columnas Mes)
    ws2 = wb["Pivot_Region_Mes"]
    headers2 = [ws2.cell(1, c).value for c in range(1, ws2.max_column + 1)]
    assert "Region" in headers2 and "Q1" in headers2 and "Q2" in headers2
    q1_idx = headers2.index("Q1") + 1
    q2_idx = headers2.index("Q2") + 1

    # Norte: Q1 = 1000 + 1500 = 2500; Q2 = 1200
    assert ws2.cell(2, 1).value == "Norte"
    assert ws2.cell(2, q1_idx).value == 2500
    assert ws2.cell(2, q2_idx).value == 1200

    # Sur: Q1 = 2000; Q2 = 800 + 2200 = 3000
    assert ws2.cell(3, 1).value == "Sur"
    assert ws2.cell(3, q1_idx).value == 2000
    assert ws2.cell(3, q2_idx).value == 3000
    assert ws2.auto_filter.ref is not None
    wb.close()

    # 4. Verificar soporte en create_excel directamente vía especificación pivot
    xlsx_direct = str(tmp_path / "pivot_direct.xlsx")
    create_excel(
        xlsx_direct,
        sheets=[
            {
                "name": "Datos",
                "headers": ["Categoria", "Unidades"],
                "rows": [["Hardware", 50], ["Software", 120], ["Hardware", 75]],
            },
            {
                "name": "Pivot_Hardware",
                "pivot": {
                    "source_sheet": "Datos",
                    "rows": "Categoria",
                    "values": "Unidades",
                    "agg": "sum",
                },
            },
        ],
    )
    wb_d = openpyxl.load_workbook(xlsx_direct)
    assert "Pivot_Hardware" in wb_d.sheetnames
    ws_d = wb_d["Pivot_Hardware"]
    assert ws_d.cell(2, 1).value == "Hardware"
    assert ws_d.cell(2, 2).value == 125
    assert ws_d.cell(3, 1).value == "Software"
    assert ws_d.cell(3, 2).value == 120
    wb_d.close()


def test_create_book_core(tmp_path):
    """Test obligatorio v0.9.0 (B): create_book multi-capítulo PDF con portada, TOC automático y EPUB opcional."""
    import fitz
    from office_worker.core.book import create_book

    out_pdf = str(tmp_path / "libro_corporativo.pdf")
    chapters = [
        {
            "title": "Introducción y Alcance",
            "content_html": "<p>Primer capítulo sobre los objetivos corporativos del año 2026.</p>",
        },
        {
            "title": "Metodología y Arquitectura",
            "content_html": "<p>Segundo capítulo explicando el diseño técnico y los pipelines documentales.</p>",
        },
        {
            "title": "Conclusiones y Próximos Pasos",
            "content_html": "<p>Tercer capítulo con los acuerdos y roadmap de ejecución.</p>",
        },
    ]

    res = create_book(
        out_path=out_pdf,
        title="Manual de Procedimientos 2026",
        author="Julio Cardozo",
        chapters=chapters,
        theme="corporate-blue",
        epub=True,
    )

    assert res["status"] == "ok"
    assert os.path.exists(res["path"])
    assert res["chapters_count"] == 3
    assert res["bytes"] > 1000

    # 1. Verificar PDF: portada, TOC, y capítulos
    doc = fitz.open(res["path"])
    assert len(doc) >= 5, f"Se esperaban al menos 5 páginas (portada + TOC + 3 caps), hay {len(doc)}"

    # Portada en página 1
    cover_text = doc[0].get_text()
    assert "Manual de Procedimientos 2026" in " ".join(cover_text.split())
    assert "Julio Cardozo" in cover_text

    # TOC en página 2
    toc_text = doc[1].get_text()
    assert "Índice General" in toc_text or "Indice General" in toc_text
    assert "Capítulo 1: Introducción y Alcance" in toc_text
    assert "Capítulo 2: Metodología y Arquitectura" in toc_text
    assert "Capítulo 3: Conclusiones y Próximos Pasos" in toc_text

    # Capítulos numerados y con contenido
    chap1_text = doc[2].get_text()
    assert "Capítulo 1" in chap1_text or "CAPÍTULO 1" in chap1_text
    assert "Introducción y Alcance" in chap1_text

    chap2_text = doc[3].get_text()
    assert "Capítulo 2" in chap2_text or "CAPÍTULO 2" in chap2_text
    assert "Metodología y Arquitectura" in chap2_text

    chap3_text = doc[4].get_text()
    assert "Capítulo 3" in chap3_text or "CAPÍTULO 3" in chap3_text
    assert "Conclusiones y Próximos Pasos" in chap3_text
    doc.close()

    # 2. Verificar exportación EPUB y magic bytes PK
    if "epub_path" in res and res["epub_path"]:
        assert os.path.exists(res["epub_path"])
        assert res["epub_bytes"] > 500
        with open(res["epub_path"], "rb") as f_epub:
            header = f_epub.read(4)
            assert header == b"PK\x03\x04" or header[:2] == b"PK", f"Magic bytes PK inválidos en EPUB: {header}"


def test_render_document_design_mode_premium(tmp_path):
    """Test obligatorio v0.9.0 (C): render_document con design_mode='premium' vs 'standard' produce bytes distintos."""
    out_std = str(tmp_path / "doc_standard.pdf")
    out_prem = str(tmp_path / "doc_premium.pdf")

    html_sample = """
    <h1>Informe Trimestral de Gestión</h1>
    <div class="kicker">Dirección General</div>
    <h2>Resumen de Operaciones</h2>
    <p>Este informe detalla el progreso y desempeño financiero alcanzado durante el período.</p>
    <div class="card"><strong>Nota:</strong> Los indicadores superaron las expectativas.</div>
    <table>
      <thead><tr><th>Indicador</th><th>Valor</th></tr></thead>
      <tbody><tr><td>Crecimiento</td><td>+18%</td></tr></tbody>
    </table>
    """

    # Modo estándar
    path_std = render_pdf(html_sample, out_std, design_mode="standard", theme="corporate-blue")
    assert os.path.exists(path_std)
    bytes_std = open(path_std, "rb").read()
    assert bytes_std[:5] == b"%PDF-"

    # Modo premium
    path_prem = render_pdf(html_sample, out_prem, design_mode="premium", theme="corporate-blue")
    assert os.path.exists(path_prem)
    bytes_prem = open(path_prem, "rb").read()
    assert bytes_prem[:5] == b"%PDF-"

    # Verificar que son distintos visualmente (bytes diferentes)
    assert bytes_std != bytes_prem, "design_mode='premium' debe producir un archivo PDF visual y binariamente distinto al estándar"
    assert len(bytes_prem) > 500
    assert len(bytes_std) > 500


def test_doctor_environment_core(tmp_path):
    """Test obligatorio v0.9.0 (D): check_environment y CLI owi doctor reportan capacidades e install_hints sin efectos secundarios."""
    import subprocess
    from office_worker.core.doctor import check_environment, detect_os, get_install_hint

    # 1. Ejecutar check_environment en core
    env_info = check_environment()
    assert env_info["status"] == "ok"
    assert "os" in env_info and isinstance(env_info["os"], str)
    assert "package_manager" in env_info and isinstance(env_info["package_manager"], str)
    assert "capabilities" in env_info and isinstance(env_info["capabilities"], dict)

    caps = env_info["capabilities"]
    required_keys = [
        "convert_to_pdf", "pdf_ocr", "poppler", "render_document",
        "pdf_core", "pdf_to_excel", "excel_core", "word_core",
        "pptx_core", "pandas_pivot", "book_epub",
    ]
    for k in required_keys:
        assert k in caps, f"Falta capability '{k}' en reporte de doctor"
        cap = caps[k]
        assert "active" in cap and isinstance(cap["active"], bool)
        assert "name" in cap and isinstance(cap["name"], str)
        assert "required_for" in cap and isinstance(cap["required_for"], str)
        assert "install_hint" in cap and isinstance(cap["install_hint"], str)

        # Si no está activo, debe proveer install_hint no vacío
        if not cap["active"]:
            assert len(cap["install_hint"]) > 0, f"Capability inactiva '{k}' debe tener un comando install_hint"

    # 2. Verificar get_install_hint determinista para Linux apt
    hint_libreoffice = get_install_hint("convert_to_pdf", "apt")
    assert "apt" in hint_libreoffice and "libreoffice" in hint_libreoffice
    hint_tesseract = get_install_hint("pdf_ocr", "apt")
    assert "apt" in hint_tesseract and "tesseract" in hint_tesseract

    # 3. CLI owi doctor ejecutado como subproceso
    cmd = [sys.executable, "-m", "office_worker.cli", "doctor"]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    assert proc.returncode == 0
    import json
    cli_data = json.loads(proc.stdout)
    assert cli_data["status"] == "ok"
    assert "capabilities" in cli_data
    assert "convert_to_pdf" in cli_data["capabilities"]







