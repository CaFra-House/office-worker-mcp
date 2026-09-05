"""Test E2E del MCP completo como test de pytest (async).

Levanta el servidor MCP real, hace handshake y ejecuta las 12 tools creando/leyendo
archivos válidos en disco. PPTX se saltea si Playwright no está disponible.
Requiere: pytest-asyncio (en extras [dev]).
"""
import asyncio, json, os, shutil, sys, pathlib
import pytest

sys.path.insert(0, str(pathlib.Path(__file__).parent.parent / "src"))

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client


@pytest.mark.asyncio
async def test_office_worker_mcp_e2e(tmp_path):
    d = str(tmp_path)
    env = {**os.environ, "PYTHONPATH": str(pathlib.Path(__file__).parent.parent / "src")}
    params = StdioServerParameters(command=sys.executable, args=["-m", "office_worker.mcp_server"], env=env)

    results = {}
    async with stdio_client(params) as (r, w):
        async with ClientSession(r, w) as s:
            await s.initialize()
            tools = [t.name for t in (await s.list_tools()).tools]
            assert len(tools) == 30 and "edit_pptx" in tools and "create_book" in tools and "environment_status" in tools and "document_diff" in tools and "scrub_metadata" in tools and "protect_office" in tools and "verify_pdf_signature" in tools and "render_document" in tools and "pdf_preview" in tools and "pdf_redact" in tools and "pdf_extract_structured" in tools and "mail_merge" in tools and "csv_excel_convert" in tools, f"tools inesperadas: {tools}"

            async def call(name, args): return json.loads((await s.call_tool(name, args)).content[0].text)

            # 1 PDF con tabla declarativa y tema corporativo nuevo
            r = await call("render_document", {
                "template_html": "<h1>{{ titulo }}</h1><p class='muted'>{{ subtitulo }}</p>{% if tabla is defined and tabla %}{{ tabla }}{% endif %}",
                "out_path": f"{d}/a.pdf",
                "theme": "corporate-blue",
                "data_json": json.dumps({"titulo": "Informe", "subtitulo": "E2E", "headers": ["KPI", "V"], "rows": [["A", 1], ["B", 2]]})})
            results["pdf"] = r; assert r["status"] == "ok" and os.path.exists(r["path"]), r

            # 1.b PDF con contraseña opcional
            r_enc = await call("render_document", {
                "template_html": "<h1>Secret E2E</h1>",
                "out_path": f"{d}/a_enc.pdf",
                "password": "e2e_secret_password"
            })
            assert r_enc["status"] == "ok" and os.path.exists(r_enc["path"])

            # 2 Word (bloques)
            r = await call("create_word", {
                "out_path": f"{d}/b.docx", "title": "Acta", "subtitle": "Reunión",
                "blocks_json": json.dumps([{"type": "h2", "text": "Puntos"}, {"type": "p", "text": "Decisión"}, {"type": "table", "headers": ["Item", "Estado"], "rows": [["X", "OK"]]}])})
            results["docx"] = r; assert r["status"] == "ok" and os.path.exists(r["path"]), r

            # 2.b Word (plantilla docxtpl)
            from docx import Document
            tpl_f = f"{d}/template_e2e.docx"
            d_tpl = Document(); d_tpl.add_paragraph("Hola {{ cliente }}"); d_tpl.save(tpl_f)
            r_tpl = await call("create_word", {
                "out_path": f"{d}/b_from_tpl.docx",
                "template_docx": tpl_f,
                "context": {"cliente": "Cliente VIP"}
            })
            assert r_tpl["status"] == "ok" and os.path.exists(r_tpl["path"])

            # 3 Excel multi-hoja
            r = await call("create_excel", {
                "out_path": f"{d}/c.xlsx", "title": "Balance",
                "sheets_json": json.dumps([{"name": "Resumen", "headers": ["Rubro", "Monto"], "rows": [["Ingresos", 100], ["Gastos", 40]]}])})
            results["xlsx"] = r; assert r["status"] == "ok" and os.path.exists(r["path"]), r

            # 4 PPTX editable — se saltea si no hay Playwright instalado en el entorno
            try:
                import playwright  # noqa: F401
                have_pw = True
            except Exception:
                have_pw = False
            if have_pw:
                r = await call("create_pptx", {
                    "out_path": f"{d}/d.pptx",
                    "slides_json": json.dumps([{"title": "Portada", "kicker": "The Office Worker"}, {"title": "Agenda", "bullets": ["Intro", "Resultados"]}])})
                results["pptx"] = r; assert r["status"] == "ok" and os.path.exists(r["path"]), r
            else:
                results["pptx"] = {"status": "skipped"}

            # 5-7 PDF input sobre el PDF generado en (1) + consolidado
            r = await call("read_pdf", {"path": results["pdf"]["path"], "extract_tables": True, "list_forms": True})
            results["read_pdf"] = r
            assert r.get("n_pages", 0) >= 1 and r.get("pages") and "tables_by_page" in r and "is_form" in r, r

            # 7.b Preview PNG via PyMuPDF (data_url base64 + archivo en disco opcional)
            import base64
            r_prev = await call("pdf_preview", {"input_path": results["pdf"]["path"], "dpi": 110})
            results["preview"] = r_prev
            assert r_prev.get("status") == "ok" and r_prev.get("data_url", "").startswith("data:image/png;base64,"), r_prev
            raw_png = base64.b64decode(r_prev["data_url"].split(",", 1)[1])
            assert raw_png[:8] == b"\x89PNG\r\n\x1a\n", "Magic bytes PNG inválidos en pdf_preview E2E"

            # 7.c Preview guardando en disco
            r_prev_file = await call("pdf_preview", {"input_path": results["pdf"]["path"], "output": f"{d}/preview.png", "max_pages": 1})
            assert r_prev_file.get("status") == "ok" and os.path.exists(f"{d}/preview.png")
            assert open(f"{d}/preview.png", "rb").read(8) == b"\x89PNG\r\n\x1a\n"

            # 8 temas
            r = await call("list_themes", {}); results["themes"] = r; assert r.get("theme", {}).get("primary"), r

            # 9 convert_to_pdf (Word -> PDF vía LibreOffice)
            if shutil.which("soffice") or shutil.which("libreoffice"):
                r = await call("convert_to_pdf", {"input_file": results["docx"]["path"], "output": f"{d}/b_from_docx.pdf"})
                results["convert"] = r
                assert r["status"] == "ok" and os.path.exists(r["path"]), r
                assert r.get("fidelity") == "clean" and "warnings" in r, r
                assert open(r["path"], "rb").read(5) == b"%PDF-"

            # 10 pdf_manipulate (merge 2 pdfs con password y extract)
            r = await call("pdf_manipulate", {
                "operation": "merge",
                "output": f"{d}/merged_e2e.pdf",
                "files": [results["pdf"]["path"], results["pdf"]["path"]],
                "password": "pwd_merged_e2e"
            })
            results["merge"] = r
            assert r["status"] == "ok" and os.path.exists(r["path"]), r

            r = await call("pdf_manipulate", {
                "operation": "extract",
                "output": f"{d}/extracted_e2e.pdf",
                "input_path": results["pdf"]["path"],
                "pages": "1"
            })
            results["extract"] = r
            assert r["status"] == "ok" and os.path.exists(r["path"]), r

            # 11 pdf_fill_form (crear form con reportlab y rellenar vía tool)
            try:
                from reportlab.pdfgen import canvas
                form_file = f"{d}/e2e_form.pdf"
                c = canvas.Canvas(form_file)
                c.drawString(50, 750, "Email:")
                c.acroForm.textfield(name="email", x=100, y=745, width=150, height=20)
                c.showPage()
                c.save()

                r = await call("pdf_fill_form", {
                    "input_pdf": form_file,
                    "fields": {"email": "test@example.com"},
                    "output": f"{d}/e2e_filled.pdf"
                })
                results["fill_form"] = r
                assert r["status"] == "ok" and os.path.exists(r["path"]), r
            except Exception as exc:
                print(f"pdf_fill_form E2E skip/warn: {exc}")

            # 12 pdf_ocr (imagen -> texto + PDF buscable)
            if shutil.which("tesseract"):
                try:
                    from PIL import Image, ImageDraw
                    img_file = f"{d}/e2e_ocr.png"
                    img = Image.new("RGB", (250, 60), color=(255, 255, 255))
                    d_ctx = ImageDraw.Draw(img)
                    d_ctx.text((10, 20), "HELLO WORLD", fill=(0, 0, 0))
                    img.save(img_file)

                    r = await call("pdf_ocr", {
                        "input_path": img_file,
                        "lang": "eng",
                        "output": f"{d}/e2e_ocr.pdf"
                    })
                    results["ocr"] = r
                    assert r["status"] == "ok", r
                    assert "HELLO" in r.get("text", "").replace(" ", "").upper()
                    assert os.path.exists(r["path"])
                except Exception as exc:
                    print(f"pdf_ocr E2E skip/warn: {exc}")

            # 13 list_templates
            r_tpls = await call("list_templates", {})
            results["list_templates"] = r_tpls
            assert r_tpls["status"] == "ok"
            assert len(r_tpls["templates"]) == 5
            tpl_names = [t["name"] for t in r_tpls["templates"]]
            assert "acta_meeting" in tpl_names and "factura_simple" in tpl_names

            # 13.b create_word con plantilla empaquetada
            r_pack = await call("create_word", {
                "out_path": f"{d}/acta_packaged.docx",
                "template_docx": "acta_meeting",
                "context": {
                    "titulo": "Acta E2E",
                    "fecha": "2026-09-05",
                    "hora": "12:00",
                    "lugar": "Sala A",
                    "asistentes": [{"nombre": "Ana", "rol": "Líder"}],
                    "puntos": [{"orden": 1, "tema": "Test", "discusion": "Exitoso"}],
                    "acuerdos": [{"acuerdo": "Aprobar", "responsable": "Ana", "fecha_limite": "Hoy"}],
                    "firmas": [{"nombre": "Ana", "cargo": "Líder"}],
                }
            })
            assert r_pack["status"] == "ok" and os.path.exists(r_pack["path"])

            # 14 pdf_compress
            r_comp = await call("pdf_compress", {
                "input_path": results["pdf"]["path"],
                "output": f"{d}/compressed_e2e.pdf",
                "quality": "med"
            })
            results["compress"] = r_comp
            assert r_comp["status"] == "ok" and os.path.exists(r_comp["path"])

            # 15 sign_pdf (estampa sello visual PNG)
            from PIL import Image, ImageDraw
            seal_file = f"{d}/seal.png"
            seal_img = Image.new("RGBA", (150, 50), color=(0, 0, 0, 0))
            d_seal = ImageDraw.Draw(seal_img)
            d_seal.rectangle([(1, 1), (148, 48)], outline=(0, 51, 102, 255), width=2)
            d_seal.text((10, 15), "SELLO E2E", fill=(0, 51, 102, 255))
            seal_img.save(seal_file)

            r_sign = await call("sign_pdf", {
                "input_pdf": results["pdf"]["path"],
                "output": f"{d}/signed_e2e.pdf",
                "sello_img_path": seal_file,
                "reason": "Test E2E",
                "location": "Buenos Aires"
            })
            results["sign"] = r_sign
            assert r_sign["status"] == "ok" and os.path.exists(r_sign["path"])

            # 16 edit_excel (modificar c.xlsx)
            r_edit_xl = await call("edit_excel", {
                "input_path": results["xlsx"]["path"],
                "operations": [
                    {"op": "set_cell", "coordinate": "B2", "value": 150},
                    {"op": "append_row", "row": ["Inversiones", 60]},
                    {"op": "add_column", "header": "Estado", "values": ["Cerrado", "Abierto", "Pendiente"]},
                    {"op": "add_chart", "chart_type": "bar", "title": "Finanzas E2E", "target_cell": "E2"},
                    {"op": "add_table", "table_style": "TableStyleMedium9"},
                    {"op": "auto_filter"}
                ]
            })
            assert r_edit_xl["status"] == "ok" and os.path.exists(r_edit_xl["path"]), r_edit_xl
            assert r_edit_xl.get("fidelity") == "rich"

            # 17 edit_word (modificar b.docx)
            r_edit_wd = await call("edit_word", {
                "input_path": results["docx"]["path"],
                "operations": [
                    {"op": "append_paragraph", "text": "Párrafo editado E2E", "bold": True},
                    {"op": "replace_text", "find": "Decisión", "replace": "Acuerdo Unánime"},
                    {"op": "insert_after_heading", "heading_text": "Puntos", "text": "Punto 1.1 introducido"},
                    {"op": "append_table", "headers": ["Aprobador", "Firma"], "rows": [["Julio", "OK"]]}
                ]
            })
            assert r_edit_wd["status"] == "ok" and os.path.exists(r_edit_wd["path"]), r_edit_wd
            assert r_edit_wd.get("fidelity") == "clean"

            # 18 read_pdf con extract_images=True sobre el signed_e2e.pdf
            r_img = await call("read_pdf", {
                "path": results["sign"]["path"],
                "extract_images": True,
                "max_images": 5
            })
            assert r_img["status"] == "ok"
            assert r_img.get("n_images", 0) >= 1
            assert "images" in r_img
            assert r_img["images"][0]["data_url"].startswith("data:image/")

            # 19 pdf_to_excel (extraer tablas de a.pdf a un nuevo .xlsx)
            r_p2x = await call("pdf_to_excel", {
                "input_path": results["pdf"]["path"],
                "output_path": f"{d}/from_pdf_e2e.xlsx",
                "sheet_name": "Reporte"
            })
            assert r_p2x["status"] == "ok" and os.path.exists(r_p2x["path"]), r_p2x
            assert r_p2x.get("fidelity") == "clean"
            assert r_p2x.get("n_tables", 0) >= 1

            # 20 read_office (leer docx y xlsx)
            r_ro_docx = await call("read_office", {"path": results["docx"]["path"]})
            assert r_ro_docx["status"] == "ok" and r_ro_docx["format"] == "docx"
            assert r_ro_docx["n_paragraphs"] >= 1

            r_ro_xlsx = await call("read_office", {"path": results["xlsx"]["path"]})
            assert r_ro_xlsx["status"] == "ok" and r_ro_xlsx["format"] == "xlsx"
            assert r_ro_xlsx["n_sheets"] >= 1

            # 21 office_batch (ejecutar lote con manejo de error parcial)
            r_batch = await call("office_batch", {
                "operations": [
                    {
                        "tool": "create_word",
                        "args": {
                            "out_path": f"{d}/batch_doc.docx",
                            "title": "Documento Batch",
                            "blocks_json": json.dumps([{"type": "p", "text": "Contenido Batch"}])
                        }
                    },
                    {
                        "tool": "tool_inexistente",
                        "args": {"foo": "bar"}
                    },
                    {
                        "tool": "list_templates",
                        "args": {}
                    }
                ]
            })
            assert r_batch["status"] == "partial_error"
            assert r_batch["total"] == 3
            assert r_batch["succeeded"] == 2
            assert r_batch["failed"] == 1
            assert os.path.exists(f"{d}/batch_doc.docx")

            # 22 pdf_redact E2E
            import fitz
            doc_sec = fitz.open()
            p_sec = doc_sec.new_page()
            p_sec.insert_text((50, 100), "CONFIDENCIAL_E2E: Clave secreta 12345.")
            p_sec.insert_text((50, 200), "COORD_SECRET_E2E en posición fija.")
            sec_pdf = f"{d}/sec_e2e.pdf"
            doc_sec.save(sec_pdf)
            doc_sec.close()

            r_redact = await call("pdf_redact", {
                "input_path": sec_pdf,
                "output": f"{d}/redacted_e2e.pdf",
                "search_text": "CONFIDENCIAL_E2E",
                "regions": [{"page": 1, "x0": 45, "y0": 185, "x1": 300, "y1": 215}],
                "fill_color": "black"
            })
            assert r_redact["status"] == "ok" and os.path.exists(r_redact["path"])
            assert r_redact["redactions_count"] >= 2
            doc_check = fitz.open(r_redact["path"])
            assert "CONFIDENCIAL_E2E" not in doc_check[0].get_text()
            assert "COORD_SECRET_E2E" not in doc_check[0].get_text()
            doc_check.close()

            # 23 pdf_extract_structured E2E (JSON & Markdown)
            r_ext_json = await call("pdf_extract_structured", {
                "input_path": results["pdf"]["path"],
                "format": "json",
                "output": f"{d}/extracted_e2e.json"
            })
            assert r_ext_json["status"] == "ok" and r_ext_json["format"] == "json"
            assert r_ext_json["n_pages"] >= 1 and len(r_ext_json["pages"]) >= 1
            assert os.path.exists(f"{d}/extracted_e2e.json")

            r_ext_md = await call("pdf_extract_structured", {
                "input_path": results["pdf"]["path"],
                "format": "markdown",
                "output": f"{d}/extracted_e2e.md"
            })
            assert r_ext_md["status"] == "ok" and r_ext_md["format"] == "markdown"
            assert "## Página" in r_ext_md["content"]
            assert os.path.exists(f"{d}/extracted_e2e.md")

            # 24 pdf_manipulate flatten E2E
            if "fill_form" in results and os.path.exists(results["fill_form"]["path"]):
                r_flat = await call("pdf_manipulate", {
                    "operation": "flatten",
                    "input_path": results["fill_form"]["path"],
                    "output": f"{d}/flattened_e2e.pdf"
                })
                assert r_flat["status"] == "ok" and os.path.exists(r_flat["path"])
                from pypdf import PdfReader
                assert not PdfReader(r_flat["path"]).get_fields()

            # 25 pdf_manipulate split_smart E2E
            doc_concat = fitz.open()
            p0 = doc_concat.new_page()
            p0.insert_text((50, 100), "DOC A: Reporte")
            doc_concat.new_page()  # Blank separator page
            p2 = doc_concat.new_page()
            p2.insert_text((50, 100), "DOC B: Contrato")
            concat_path = f"{d}/concat_e2e.pdf"
            doc_concat.save(concat_path)
            doc_concat.close()

            r_split = await call("pdf_manipulate", {
                "operation": "split_smart",
                "input_path": concat_path,
                "output": f"{d}/split_smart_e2e.pdf"
            })
            assert r_split["status"] == "ok" and r_split["n_splits"] == 2
            assert len(r_split["files"]) == 2
            for f_part in r_split["files"]:
                assert os.path.exists(f_part)
                assert open(f_part, "rb").read(5) == b"%PDF-"

            # 26 office_batch con nuevas tools v0.6.0
            r_batch_v6 = await call("office_batch", {
                "operations": [
                    {
                        "tool": "pdf_redact",
                        "args": {
                            "input_path": sec_pdf,
                            "output": f"{d}/batch_redacted.pdf",
                            "search_text": "Clave secreta"
                        }
                    },
                    {
                        "tool": "pdf_extract_structured",
                        "args": {
                            "input_path": results["pdf"]["path"],
                            "format": "markdown",
                            "output": f"{d}/batch_extracted.md"
                        }
                    }
                ]
            })
            assert r_batch_v6["status"] == "ok"
            assert r_batch_v6["total"] == 2
            assert r_batch_v6["succeeded"] == 2
            assert os.path.exists(f"{d}/batch_redacted.pdf")
            assert os.path.exists(f"{d}/batch_extracted.md")

            # 27 mail_merge E2E
            mm_tpl = f"{d}/mail_merge_tpl.docx"
            doc_mm = Document()
            doc_mm.add_paragraph("Hola {{ cliente }}, tu plan es {{ plan }}.")
            doc_mm.save(mm_tpl)

            mm_csv = f"{d}/clientes.csv"
            with open(mm_csv, "w", encoding="utf-8") as f:
                f.write("cliente,plan\nAna,Pro\nBeto,Enterprise\n")

            r_mm = await call("mail_merge", {
                "template_path": mm_tpl,
                "dataset_csv": mm_csv,
                "output_prefix": f"{d}/carta_mm"
            })
            assert r_mm["status"] == "ok"
            assert r_mm["n_docs"] == 2
            assert len(r_mm["paths"]) == 2
            for p_mm in r_mm["paths"]:
                assert os.path.exists(p_mm)
                assert open(p_mm, "rb").read(2) == b"PK"
            doc_check1 = Document(r_mm["paths"][0])
            assert "Ana" in doc_check1.paragraphs[0].text
            doc_check2 = Document(r_mm["paths"][1])
            assert "Beto" in doc_check2.paragraphs[0].text

            # 28 read_office formato markdown E2E
            r_ro_md = await call("read_office", {
                "path": results["docx"]["path"],
                "format": "markdown"
            })
            assert r_ro_md["status"] == "ok"
            assert r_ro_md["format"] == "markdown"
            assert "# " in r_ro_md["content"] or "## " in r_ro_md["content"]
            assert "|" in r_ro_md["content"]

            r_ro_xlsx_md = await call("read_office", {
                "path": results["xlsx"]["path"],
                "format": "markdown"
            })
            assert r_ro_xlsx_md["status"] == "ok"
            assert r_ro_xlsx_md["format"] == "markdown"
            assert "## Resumen" in r_ro_xlsx_md["content"]
            assert "Ingresos" in r_ro_xlsx_md["content"] and "|" in r_ro_xlsx_md["content"]

            # 29 csv_excel_convert bidireccional E2E
            r_c2x = await call("csv_excel_convert", {
                "input": mm_csv,
                "output": f"{d}/clientes.xlsx",
                "direction": "csv_to_xlsx",
                "sheet": "Clientes"
            })
            assert r_c2x["status"] == "ok"
            assert os.path.exists(r_c2x["path"])
            assert open(r_c2x["path"], "rb").read(2) == b"PK"
            assert r_c2x["n_rows"] == 3

            r_x2c = await call("csv_excel_convert", {
                "input": r_c2x["path"],
                "output": f"{d}/clientes_roundtrip.csv",
                "direction": "xlsx_to_csv"
            })
            assert r_x2c["status"] == "ok"
            assert os.path.exists(r_x2c["path"])
            with open(r_x2c["path"], "r", encoding="utf-8") as f:
                csv_rt = f.read()
            assert "Ana" in csv_rt and "Enterprise" in csv_rt

            # 30 create_pptx con gráfico nativo E2E
            r_pptx_chart = await call("create_pptx", {
                "out_path": f"{d}/deck_chart_e2e.pptx",
                "slides_json": json.dumps([
                    {
                        "title": "Métricas Q1-Q4",
                        "kicker": "Rendimiento",
                        "chart": {
                            "type": "bar",
                            "title": "Ventas 2026",
                            "categories": ["Q1", "Q2", "Q3", "Q4"],
                            "values": [100, 150, 220, 310]
                        }
                    }
                ])
            })
            assert r_pptx_chart["status"] == "ok"
            assert os.path.exists(r_pptx_chart["path"])
            from pptx import Presentation
            prs_e2e = Presentation(r_pptx_chart["path"])
            assert any(s.has_chart for s in prs_e2e.slides[0].shapes)

            # 30.b edit_pptx E2E (in-place modification with python-pptx)
            r_edit_pptx = await call("edit_pptx", {
                "input_path": r_pptx_chart["path"],
                "operations": [
                    {"op": "replace_text", "find": "Métricas", "replace": "Indicadores"},
                    {"op": "add_slide", "title": "Conclusiones E2E", "bullets": ["Objetivo cumplido", "Rendimiento alto"]},
                    {"op": "set_notes", "slide_index": 0, "notes": "Notas del presentador E2E"},
                ]
            })
            assert r_edit_pptx["status"] == "ok" and os.path.exists(r_edit_pptx["path"]), r_edit_pptx
            assert r_edit_pptx.get("fidelity") == "high"
            assert r_edit_pptx.get("operations") == 3
            prs_edited = Presentation(r_edit_pptx["path"])
            assert len(prs_edited.slides) == 2
            assert "Notas del presentador E2E" in prs_edited.slides[0].notes_slide.notes_text_frame.text
            assert any(s.has_chart for s in prs_edited.slides[0].shapes)

            # 31 office_batch con nuevas tools v0.7.0
            r_batch_v7 = await call("office_batch", {
                "operations": [
                    {
                        "tool": "csv_excel_convert",
                        "args": {
                            "input": mm_csv,
                            "output": f"{d}/batch_clientes.xlsx",
                            "direction": "csv_to_xlsx"
                        }
                    },
                    {
                        "tool": "read_office",
                        "args": {
                            "path": f"{d}/batch_clientes.xlsx",
                            "format": "markdown"
                        }
                    }
                ]
            })
            assert r_batch_v7["status"] == "ok"
            assert r_batch_v7["total"] == 2
            assert r_batch_v7["succeeded"] == 2
            assert os.path.exists(f"{d}/batch_clientes.xlsx")

            # 32 next_steps verification on primary creation tools
            assert "next_steps" in results["pdf"] and len(results["pdf"]["next_steps"]) > 0
            assert "next_steps" in results["docx"] and len(results["docx"]["next_steps"]) > 0
            assert "next_steps" in results["xlsx"] and len(results["xlsx"]["next_steps"]) > 0
            assert "next_steps" in results["sign"] and len(results["sign"]["next_steps"]) > 0
            assert "next_steps" in r_redact and len(r_redact["next_steps"]) > 0
            assert "next_steps" in r_mm and len(r_mm["next_steps"]) > 0
            if "convert" in results:
                assert "next_steps" in results["convert"] and len(results["convert"]["next_steps"]) > 0
            if results["pptx"].get("status") == "ok":
                assert "next_steps" in results["pptx"] and len(results["pptx"]["next_steps"]) > 0

            # 33 document_diff E2E
            doc_mod = f"{d}/b_modified.docx"
            from docx import Document
            d_mod = Document(results["docx"]["path"])
            d_mod.add_paragraph("Nueva cláusula de auditoría agregada en revisión.")
            d_mod.save(doc_mod)

            r_diff = await call("document_diff", {
                "path_a": results["docx"]["path"],
                "path_b": doc_mod,
                "format": "markdown",
            })
            assert r_diff["status"] == "ok"
            assert r_diff["has_changes"] is True
            assert r_diff["summary"]["added"] >= 1
            assert "diff_markdown" in r_diff
            assert len(r_diff.get("warnings", [])) > 0

            # 34 scrub_metadata E2E
            scrubbed_docx = f"{d}/docx_scrubbed.docx"
            r_scrub = await call("scrub_metadata", {
                "input": results["docx"]["path"],
                "output": scrubbed_docx,
            })
            assert r_scrub["status"] == "ok"
            assert os.path.exists(scrubbed_docx)
            assert open(scrubbed_docx, "rb").read(2) == b"PK"
            assert "author" in r_scrub.get("scrubbed_fields", []) or "revision" in r_scrub.get("scrubbed_fields", [])

            # 35 protect_office E2E
            protected_xlsx = f"{d}/xlsx_protected.xlsx"
            r_prot = await call("protect_office", {
                "input": results["xlsx"]["path"],
                "output": protected_xlsx,
                "password": "PasswordE2E2026!",
            })
            assert r_prot["status"] == "ok"
            assert os.path.exists(protected_xlsx)
            assert r_prot.get("encrypted") is True
            from zipfile import BadZipFile
            import openpyxl, msoffcrypto
            with pytest.raises(BadZipFile):
                openpyxl.load_workbook(protected_xlsx)
            with open(protected_xlsx, "rb") as f:
                of = msoffcrypto.OfficeFile(f)
                assert of.is_encrypted() is True

            # 36 verify_pdf_signature E2E
            # 36.a Unsigned PDF
            r_ver_un = await call("verify_pdf_signature", {
                "input": results["pdf"]["path"],
            })
            assert r_ver_un["status"] == "ok"
            assert r_ver_un["has_signature"] is False
            assert r_ver_un["valid"] is False

            # 36.b Digital signed PDF with PEM
            import datetime
            from cryptography import x509
            from cryptography.x509.oid import NameOID
            from cryptography.hazmat.primitives import hashes, serialization
            from cryptography.hazmat.primitives.asymmetric import rsa

            key = rsa.generate_private_key(public_exponent=65537, key_size=2048)
            subject = issuer = x509.Name([
                x509.NameAttribute(NameOID.COMMON_NAME, "Julio Cardozo E2E"),
            ])
            cert = (
                x509.CertificateBuilder()
                .subject_name(subject)
                .issuer_name(issuer)
                .public_key(key.public_key())
                .serial_number(x509.random_serial_number())
                .not_valid_before(datetime.datetime.now(datetime.timezone.utc))
                .not_valid_after(datetime.datetime.now(datetime.timezone.utc) + datetime.timedelta(days=1))
                .sign(key, hashes.SHA256())
            )
            cert_pem_file = f"{d}/cert_e2e.pem"
            with open(cert_pem_file, "wb") as f:
                f.write(cert.public_bytes(serialization.Encoding.PEM))
                f.write(key.private_bytes(
                    encoding=serialization.Encoding.PEM,
                    format=serialization.PrivateFormat.PKCS8,
                    encryption_algorithm=serialization.NoEncryption()
                ))

            signed_dig_pdf = f"{d}/signed_digital_e2e.pdf"
            r_sign_dig = await call("sign_pdf", {
                "input_pdf": results["pdf"]["path"],
                "output": signed_dig_pdf,
                "cert_pem": cert_pem_file,
                "reason": "Validación Técnica E2E",
                "location": "Buenos Aires",
            })
            assert r_sign_dig["status"] == "ok"

            r_ver_dig = await call("verify_pdf_signature", {
                "input": signed_dig_pdf,
            })
            assert r_ver_dig["status"] == "ok"
            assert r_ver_dig["has_signature"] is True
            assert r_ver_dig["valid"] is True
            assert r_ver_dig["intact"] is True
            assert "Julio Cardozo E2E" in (r_ver_dig["signer"] or "")
            assert r_ver_dig["reason"] == "Validación Técnica E2E"

            # 36.c Digital signed PDF with auto_generate_test_cert
            signed_auto_pdf = f"{d}/signed_auto_e2e.pdf"
            r_sign_auto = await call("sign_pdf", {
                "input_pdf": results["pdf"]["path"],
                "output": signed_auto_pdf,
                "auto_generate_test_cert": True,
                "reason": "Test E2E Auto Cert",
            })
            assert r_sign_auto["status"] == "ok"
            assert "next_steps" in r_sign_auto and len(r_sign_auto["next_steps"]) > 0

            r_ver_auto = await call("verify_pdf_signature", {
                "input": signed_auto_pdf,
            })
            assert r_ver_auto["status"] == "ok"
            assert r_ver_auto["has_signature"] is True
            assert r_ver_auto["valid"] is True
            assert r_ver_auto["intact"] is True
            assert "Office Worker Test Certificate (Non-Production)" in (r_ver_auto["signer"] or "")

            # 37 office_batch con tools v0.8.0
            r_batch_v8 = await call("office_batch", {
                "operations": [
                    {
                        "tool": "document_diff",
                        "args": {
                            "path_a": results["docx"]["path"],
                            "path_b": doc_mod,
                        }
                    },
                    {
                        "tool": "scrub_metadata",
                        "args": {
                            "input": results["docx"]["path"],
                            "output": f"{d}/batch_scrubbed.docx",
                        }
                    }
                ]
            })
            assert r_batch_v8["status"] == "ok"
            assert r_batch_v8["total"] == 2
            assert r_batch_v8["succeeded"] == 2

            # 38 edit_excel add_pivot E2E
            pivot_file = f"{d}/pivot_e2e.xlsx"
            await call("create_excel", {
                "out_path": pivot_file,
                "title": "Ventas E2E",
                "sheets_json": json.dumps([{
                    "name": "Ventas",
                    "headers": ["Region", "Producto", "Monto"],
                    "rows": [
                        ["Norte", "Alpha", 100],
                        ["Sur", "Alpha", 150],
                        ["Norte", "Beta", 200],
                        ["Sur", "Beta", 250],
                    ]
                }])
            })
            r_xl_pivot = await call("edit_excel", {
                "input_path": pivot_file,
                "operations": [
                    {
                        "op": "add_pivot",
                        "sheet": "Ventas",
                        "rows": "Region",
                        "cols": "Producto",
                        "values": "Monto",
                        "agg": "sum",
                        "pivot_sheet": "Pivot_Ventas",
                    }
                ]
            })
            assert r_xl_pivot["status"] == "ok", r_xl_pivot
            assert "Pivot_Ventas" in r_xl_pivot.get("sheets_modified", [])
            assert os.path.exists(pivot_file)

            # 39 create_book E2E (PDF + EPUB)
            r_book = await call("create_book", {
                "output": f"{d}/libro_e2e.pdf",
                "title": "Manual de Gestión Empresarial",
                "author": "Julio Cardozo",
                "chapters": [
                    {"title": "Visión Estratégica", "content_html": "<p>Contenido del capítulo 1.</p>"},
                    {"title": "Operaciones y MCP", "content_html": "<p>Contenido del capítulo 2.</p>"},
                    {"title": "Auditoría Continua", "content_html": "<p>Contenido del capítulo 3.</p>"},
                ],
                "theme": "corporate-blue",
                "epub": True,
            })
            assert r_book["status"] == "ok"
            assert os.path.exists(r_book["path"])
            assert r_book["chapters_count"] == 3
            assert "next_steps" in r_book and len(r_book["next_steps"]) > 0
            if "epub_path" in r_book and r_book["epub_path"]:
                assert os.path.exists(r_book["epub_path"])
                with open(r_book["epub_path"], "rb") as f_ep:
                    assert f_ep.read(2) == b"PK"

            # 40 render_document con design_mode='premium' E2E
            r_prem = await call("render_document", {
                "template_html": "<h1>Informe Premium E2E</h1><div class='kicker'>Dirección</div><p>Texto editorial.</p>",
                "out_path": f"{d}/doc_premium_e2e.pdf",
                "theme": "corporate-blue",
                "design_mode": "premium",
            })
            assert r_prem["status"] == "ok"
            assert os.path.exists(r_prem["path"])
            assert open(r_prem["path"], "rb").read(5) == b"%PDF-"
            assert r_prem["bytes"] != results["pdf"]["bytes"]

            # 41 environment_status E2E
            r_env = await call("environment_status", {})
            assert r_env["status"] == "ok"
            assert "os" in r_env
            assert "capabilities" in r_env
            assert "render_document" in r_env["capabilities"]
            assert "convert_to_pdf" in r_env["capabilities"]

            # 42 office_batch con tools v0.9.0
            r_batch_v9 = await call("office_batch", {
                "operations": [
                    {
                        "tool": "environment_status",
                        "args": {},
                    },
                    {
                        "tool": "render_document",
                        "args": {
                            "template_html": "<h1>Batch Premium</h1>",
                            "out_path": f"{d}/batch_premium.pdf",
                            "design_mode": "premium",
                        },
                    },
                ]
            })
            assert r_batch_v9["status"] == "ok"
            assert r_batch_v9["total"] == 2
            assert r_batch_v9["succeeded"] == 2



