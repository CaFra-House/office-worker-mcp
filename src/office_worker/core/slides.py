"""Generación de PPTX editable vía html-to-pptx (Playwright Chromium) o nativo python-pptx con soporte de gráficos."""
from __future__ import annotations
import os, tempfile
from typing import Any

_SLIDE_CSS = """
* { box-sizing: border-box; margin: 0; padding: 0; }
body { font-family: var(--ow-font-body); color: var(--ow-text); background: #f1f5f9; }
.slide {
  width: 1280px;
  height: 720px;
  padding: 48px 64px;
  background: var(--ow-bg);
  display: flex;
  flex-direction: column;
  position: relative;
  overflow: hidden;
}
.slide.cover {
  justify-content: center;
  padding: 72px 88px;
  border-top: 8px solid var(--ow-primary);
}
.slide.cover .cover-accent-bar {
  width: 6px;
  height: 100%;
  position: absolute;
  left: 0;
  top: 0;
  background: var(--ow-accent);
}
.slide.cover .kicker {
  color: var(--ow-accent);
  font-weight: 700;
  letter-spacing: 0.15em;
  text-transform: uppercase;
  font-size: 16px;
  margin-bottom: 14px;
}
.slide.cover h1 {
  font-family: var(--ow-font-title);
  font-size: 48px;
  font-weight: 800;
  color: var(--ow-primary);
  line-height: 1.15;
  border-bottom: none;
  padding-bottom: 0;
  margin-bottom: 16px;
}
.slide.cover .subtitle {
  font-size: 24px;
  color: var(--ow-muted);
  line-height: 1.4;
  margin-top: 6px;
}
.slide.cover ul {
  margin-top: 20px;
  list-style: none;
  padding-left: 0;
}
.slide.cover li {
  font-size: 18px;
  color: var(--ow-text);
  margin-bottom: 8px;
}
.slide.cover li::before {
  content: "▪ ";
  color: var(--ow-accent);
  font-weight: bold;
}
.slide .header-block {
  margin-bottom: 24px;
  border-bottom: 2.5px solid var(--ow-accent);
  padding-bottom: 12px;
}
.slide .header-block .kicker {
  color: var(--ow-accent);
  font-weight: 700;
  letter-spacing: 0.12em;
  text-transform: uppercase;
  font-size: 13px;
  margin-bottom: 4px;
}
.slide .header-block h1 {
  font-family: var(--ow-font-title);
  font-size: 34px;
  font-weight: 700;
  color: var(--ow-primary);
  border-bottom: none;
  padding-bottom: 0;
  margin: 0;
}
.slide h2 {
  font-family: var(--ow-font-title);
  font-size: 26px;
  color: var(--ow-primary);
  margin-bottom: 16px;
}
.slide p {
  font-size: 20px;
  line-height: 1.5;
  color: var(--ow-text);
}
.slide .subtitle {
  font-size: 20px;
  color: var(--ow-muted);
  font-style: italic;
  margin-bottom: 16px;
}
.slide ul {
  padding-left: 0;
  margin-top: 8px;
  list-style: none;
}
.slide li {
  font-size: 19px;
  line-height: 1.6;
  margin-bottom: 12px;
  color: var(--ow-text);
  padding-left: 24px;
  position: relative;
}
.slide li::before {
  content: "▪";
  color: var(--ow-accent);
  font-weight: bold;
  position: absolute;
  left: 0;
  top: 0;
}
.slide table {
  width: 100%;
  border-collapse: separate;
  border-spacing: 0;
  font-size: 16px;
  margin-top: 16px;
  border: 1px solid #cbd5e1;
  border-radius: 6px;
  overflow: hidden;
}
.slide th {
  background: var(--ow-primary);
  color: #ffffff;
  font-family: var(--ow-font-title);
  font-weight: 700;
  text-transform: uppercase;
  letter-spacing: 0.08em;
  font-size: 14px;
  padding: 12px 16px;
  text-align: left;
}
.slide td {
  padding: 10px 16px;
  border-bottom: 1px solid #e2e8f0;
  color: var(--ow-text);
  vertical-align: middle;
}
.slide tr:last-child td {
  border-bottom: none;
}
.slide tr:nth-child(even) td {
  background: var(--ow-row-alt);
}
.slide td:first-child {
  font-weight: 600;
}
"""


def _hex_to_rgb(h: str):
    from pptx.dml.color import RGBColor
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _clean_font_name(font_str: str | None, default: str = "Arial") -> str:
    if not font_str:
        return default
    first = str(font_str).split(",")[0].strip().strip("'").strip('"')
    return first or default


def _add_chart_to_slide(slide, chart_spec: dict[str, Any]):
    """Añade un gráfico nativo de PowerPoint a la diapositiva usando python-pptx."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    c_type = str(chart_spec.get("type") or chart_spec.get("chart_type") or "bar").lower().strip()
    if c_type in ("pie", "piechart", "pie_chart"):
        xl_type = XL_CHART_TYPE.PIE
    elif c_type in ("line", "linechart", "line_chart"):
        xl_type = XL_CHART_TYPE.LINE
    elif c_type in ("horizontal_bar", "bar_horizontal"):
        xl_type = XL_CHART_TYPE.BAR_CLUSTERED
    else:
        xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    cdata = CategoryChartData()
    categories = chart_spec.get("categories") or []
    cdata.categories = [str(c) for c in categories]

    if "series" in chart_spec and isinstance(chart_spec["series"], list):
        for s in chart_spec["series"]:
            name = str(s.get("name", "Serie"))
            vals = tuple(s.get("values", []))
            cdata.add_series(name, vals)
    else:
        vals = tuple(chart_spec.get("values") or [])
        s_name = str(chart_spec.get("series_name") or chart_spec.get("title") or "Valores")
        cdata.add_series(s_name, vals)

    left = Inches(float(chart_spec.get("left", 1.2)))
    top = Inches(float(chart_spec.get("top", 1.8)))
    width = Inches(float(chart_spec.get("width", 10.5)))
    height = Inches(float(chart_spec.get("height", 4.8)))

    chart_shape = slide.shapes.add_chart(xl_type, left, top, width, height, cdata)
    chart = chart_shape.chart
    if chart_spec.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = str(chart_spec["title"])
    return chart_shape


def _create_native_pptx(out_path: str, slides: list[dict], theme: dict) -> str:
    """Generador PPTX nativo de calidad corporativo-agencia usando python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    primary_color = _hex_to_rgb(theme.get("primary", "#003366"))
    accent_color = _hex_to_rgb(theme.get("accent", "#3B82F6"))
    text_color = _hex_to_rgb(theme.get("text", "#1A202C"))
    muted_color = _hex_to_rgb(theme.get("muted", "#718096"))
    row_alt_color = _hex_to_rgb(theme.get("row_alt", "#F5F7FA"))

    font_title = _clean_font_name(theme.get("font_title"), "Helvetica")
    font_body = _clean_font_name(theme.get("font_body"), "Arial")

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)

        title_text = s.get("title", "")
        kicker_text = s.get("kicker", "")
        subtitle_text = s.get("subtitle", "")
        bullets = s.get("bullets", [])
        has_chart = "chart" in s and s["chart"]
        has_table = ("table" in s and s["table"]) or ("headers" in s and "rows" in s)

        # Portada: cuando se especifica explícitamente o en slide 0 sin chart ni tabla
        is_cover = (
            s.get("is_cover") is True
            or str(s.get("type", "")).lower() == "cover"
            or (
                idx == 0
                and s.get("is_cover") is not False
                and str(s.get("type", "")).lower() not in ("slide", "content", "internal")
                and not has_chart
                and not has_table
            )
        )

        if is_cover:
            # === PORTADA CORPORATIVO-AGENCIA ===
            # Banda de color primary arriba
            top_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.14))
            top_stripe.fill.solid()
            top_stripe.fill.fore_color.rgb = primary_color
            top_stripe.line.fill.background()

            # Barra vertical acento
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(0.08), Inches(2.6))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_color
            bar.line.fill.background()

            # Caja de texto principal con aire generoso
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.9), Inches(10.5), Inches(4.8))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_bottom = Inches(0)

            first_p = True
            if kicker_text:
                pk = tf.paragraphs[0]
                pk.text = str(kicker_text).upper()
                pk.font.size = Pt(14)
                pk.font.bold = True
                pk.font.color.rgb = accent_color
                pk.font.name = font_title
                pk.space_after = Pt(12)
                first_p = False

            if title_text:
                pt = tf.paragraphs[0] if first_p else tf.add_paragraph()
                pt.text = str(title_text)
                pt.font.size = Pt(38)
                pt.font.bold = True
                pt.font.color.rgb = primary_color
                pt.font.name = font_title
                pt.space_after = Pt(14)
                first_p = False

            if subtitle_text:
                ps = tf.paragraphs[0] if first_p else tf.add_paragraph()
                ps.text = str(subtitle_text)
                ps.font.size = Pt(20)
                ps.font.color.rgb = muted_color
                ps.font.name = font_body
                ps.space_after = Pt(12)
                first_p = False

            if bullets:
                for b_item in bullets:
                    pb = tf.paragraphs[0] if first_p else tf.add_paragraph()
                    r_sym = pb.add_run()
                    r_sym.text = "▪  "
                    r_sym.font.bold = True
                    r_sym.font.size = Pt(14)
                    r_sym.font.color.rgb = accent_color
                    r_txt = pb.add_run()
                    r_txt.text = str(b_item)
                    r_txt.font.size = Pt(16)
                    r_txt.font.color.rgb = text_color
                    r_txt.font.name = font_body
                    pb.space_before = Pt(4)
                    pb.space_after = Pt(6)
                    first_p = False

        else:
            # === SLIDE INTERNO ===
            # Header box (kicker + title)
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.55), Inches(11.333), Inches(1.2))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_bottom = Inches(0)

            first_p = True
            if kicker_text:
                p_k = tf.paragraphs[0]
                p_k.text = str(kicker_text).upper()
                p_k.font.size = Pt(12)
                p_k.font.bold = True
                p_k.font.color.rgb = accent_color
                p_k.font.name = font_title
                p_k.space_after = Pt(4)
                first_p = False

            if title_text:
                p_t = tf.paragraphs[0] if first_p else tf.add_paragraph()
                p_t.text = str(title_text)
                p_t.font.size = Pt(30)
                p_t.font.bold = True
                p_t.font.color.rgb = primary_color
                p_t.font.name = font_title

            # Línea inferior fina accent debajo del título
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.82), Inches(11.333), Inches(0.025))
            sep.fill.solid()
            sep.fill.fore_color.rgb = accent_color
            sep.line.fill.background()

            # Contenido: Tablas / Bullets / Subtitle / Chart
            content_top = Inches(2.15)
            content_height = Inches(4.7)

            # 1. Tabla si existe
            if has_table:
                tbl_spec = s.get("table") or {}
                headers = tbl_spec.get("headers") or s.get("headers") or []
                rows = tbl_spec.get("rows") or s.get("rows") or []
                n_rows = len(rows) + (1 if headers else 0)
                n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
                if n_rows > 0 and n_cols > 0:
                    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(1.0), content_top, Inches(11.333), min(Inches(0.55 * n_rows), content_height))
                    tbl = tbl_shape.table
                    # Header
                    if headers:
                        for c_i, hv in enumerate(headers):
                            c = tbl.cell(0, c_i)
                            c.fill.solid()
                            c.fill.fore_color.rgb = primary_color
                            c.text = str(hv).upper()
                            for p in c.text_frame.paragraphs:
                                for r in p.runs:
                                    r.font.bold = True
                                    r.font.size = Pt(12.5)
                                    r.font.color.rgb = RGBColor(255, 255, 255)
                                    r.font.name = font_title
                    # Filas de datos
                    start_r = 1 if headers else 0
                    for r_i, row_data in enumerate(rows):
                        is_alt = (r_i % 2 == 1)
                        for c_i, cv in enumerate(row_data):
                            if c_i < n_cols:
                                c = tbl.cell(r_i + start_r, c_i)
                                if is_alt:
                                    c.fill.solid()
                                    c.fill.fore_color.rgb = row_alt_color
                                c.text = str(cv)
                                for p in c.text_frame.paragraphs:
                                    for r in p.runs:
                                        r.font.size = Pt(13)
                                        r.font.color.rgb = text_color
                                        r.font.name = font_body
                                        if c_i == 0:
                                            r.font.bold = True

            # 2. Bullets / Subtitle (cuando no hay tabla)
            elif bullets or subtitle_text:
                box_width = Inches(5.0) if has_chart else Inches(11.333)
                b_tb = slide.shapes.add_textbox(Inches(1.0), content_top, box_width, content_height)
                b_tf = b_tb.text_frame
                b_tf.word_wrap = True
                b_tf.margin_left = Inches(0)
                b_tf.margin_top = Inches(0)

                first_bp = True
                if subtitle_text:
                    sp = b_tf.paragraphs[0]
                    sp.text = subtitle_text
                    sp.font.size = Pt(18)
                    sp.font.italic = True
                    sp.font.color.rgb = muted_color
                    sp.font.name = font_body
                    sp.space_after = Pt(12)
                    first_bp = False

                for b_idx, bullet in enumerate(bullets):
                    bp = b_tf.paragraphs[0] if (first_bp and b_idx == 0) else b_tf.add_paragraph()
                    r_sym = bp.add_run()
                    r_sym.text = "▪  "
                    r_sym.font.bold = True
                    r_sym.font.size = Pt(15)
                    r_sym.font.color.rgb = accent_color
                    r_txt = bp.add_run()
                    r_txt.text = str(bullet)
                    r_txt.font.size = Pt(18)
                    r_txt.font.color.rgb = text_color
                    r_txt.font.name = font_body
                    bp.space_before = Pt(6)
                    bp.space_after = Pt(10)
                    bp.line_spacing = 1.25

            # 3. Gráficos
            if has_chart:
                chart_spec = dict(s["chart"])
                if bullets and "left" not in chart_spec:
                    chart_spec["left"] = 6.2
                    chart_spec["width"] = 6.1
                    chart_spec["top"] = 2.15
                    chart_spec["height"] = 4.7
                elif not bullets and "left" not in chart_spec:
                    chart_spec["left"] = 1.0
                    chart_spec["width"] = 11.333
                    chart_spec["top"] = 2.15
                    chart_spec["height"] = 4.7
                _add_chart_to_slide(slide, chart_spec)

    prs.save(out_path)
    return out_path


def create_pptx(out_path, slides=None, theme=None, prefer_native: bool = False):
    """slides: lista de dicts {"title":..., "kicker":..., "bullets":[...], "chart":{type, categories, values}} o {"html":"..."} crudo.

    Si slides es None → una portada genérica. Devuelve ruta absoluta del .pptx editable.
    """
    from .themes import load_theme, css_vars
    from .security import safe_out
    from pptx import Presentation

    th = load_theme(theme)
    out_path = safe_out(out_path)

    slide_list = slides or [{"title": "The Office Worker", "bullets": ["Presentación generada"]}]

    # Si se prefiere nativo directamente, saltar html_to_pptx
    if prefer_native:
        result = _create_native_pptx(out_path, slide_list, th)
        if not os.path.exists(result) or os.path.getsize(result) < 3000:
            raise RuntimeError(f"PPTX generado inválido o vacío ({result})")
        return os.path.abspath(result)

    # Intentar renderizar vía html-to-pptx si está disponible
    try:
        from html_to_pptx import convert
    except ImportError:
        convert = None

    if convert is not None:
        parts = []
        for idx, s in enumerate(slide_list):
            if "html" in s and s["html"]:
                parts.append(f'<section class="slide">{s["html"]}</section>')
                continue

            title_text = s.get("title", "")
            kicker_text = s.get("kicker", "")
            subtitle_text = s.get("subtitle", "")
            bullets = s.get("bullets", [])
            has_chart = "chart" in s and s["chart"]
            has_table = ("table" in s and s["table"]) or ("headers" in s and "rows" in s)

            is_cover = (
                s.get("is_cover") is True
                or str(s.get("type", "")).lower() == "cover"
                or (
                    idx == 0
                    and s.get("is_cover") is not False
                    and str(s.get("type", "")).lower() not in ("slide", "content", "internal")
                    and not has_chart
                    and not has_table
                )
            )

            if is_cover:
                k = f'<div class="kicker">{kicker_text}</div>' if kicker_text else ""
                t = f'<h1>{title_text}</h1>' if title_text else ""
                sub = f'<div class="subtitle">{subtitle_text}</div>' if subtitle_text else ""
                items = "".join(f"<li>{b}</li>" for b in bullets)
                b_html = f"<ul>{items}</ul>" if items else ""
                parts.append(
                    f'<section class="slide cover">'
                    f'<div class="cover-accent-bar"></div>'
                    f'{k}{t}{sub}{b_html}'
                    f'</section>'
                )
            else:
                k = f'<div class="kicker">{kicker_text}</div>' if kicker_text else ""
                t = f'<h1>{title_text}</h1>' if title_text else ""
                hdr = f'<div class="header-block">{k}{t}</div>'

                tbl_html = ""
                if has_table:
                    tbl_spec = s.get("table") or {}
                    headers = tbl_spec.get("headers") or s.get("headers") or []
                    rows = tbl_spec.get("rows") or s.get("rows") or []
                    th_str = "".join(f"<th>{h}</th>" for h in headers)
                    tr_str = "".join("<tr>" + "".join(f"<td>{c}</td>" for c in r) + "</tr>" for r in rows)
                    tbl_html = f"<table><thead><tr>{th_str}</tr></thead><tbody>{tr_str}</tbody></table>"

                sub = f'<div class="subtitle">{subtitle_text}</div>' if subtitle_text else ""
                items = "".join(f"<li>{b}</li>" for b in bullets)
                b_html = f'<ul style="{"width:50%;" if has_chart else ""}">{items}</ul>' if items else ""

                parts.append(
                    f'<section class="slide">'
                    f'{hdr}{sub}{tbl_html}{b_html}'
                    f'</section>'
                )

        deck_html = (f'<!DOCTYPE html><html><head><meta charset="utf-8">'
                     f'<style>{css_vars(th)}{_SLIDE_CSS}</style></head><body>{"".join(parts)}</body></html>')

        with tempfile.NamedTemporaryFile("w", suffix=".html", delete=False, encoding="utf-8") as tf:
            tf.write(deck_html)
            tmp_html = tf.name

        try:
            import asyncio, concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as ex:
                result = ex.submit(asyncio.run, convert(tmp_html, out_path)).result(timeout=300)

            # Inyectar gráficos nativos en diapositivas si se declararon
            has_charts = any(isinstance(s, dict) and "chart" in s and s["chart"] for s in slide_list)
            if has_charts:
                prs = Presentation(result)
                for idx, s in enumerate(slide_list):
                    if isinstance(s, dict) and s.get("chart") and idx < len(prs.slides):
                        _add_chart_to_slide(prs.slides[idx], s["chart"])
                prs.save(result)

        except Exception:
            # Fallback a generación nativa python-pptx si falló el navegador
            result = _create_native_pptx(out_path, slide_list, th)
        finally:
            try: os.unlink(tmp_html)
            except OSError: pass
    else:
        result = _create_native_pptx(out_path, slide_list, th)

    if not os.path.exists(result) or os.path.getsize(result) < 3000:
        raise RuntimeError(f"PPTX generado inválido o vacío ({result})")
    return os.path.abspath(result)


def _iter_all_shapes(shapes):
    """Recursively yields all shapes, including shapes inside group shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        yield shape
        if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)


def _replace_in_text_frame(tf, old: str, new: str, count_left: int) -> tuple[int, int]:
    """Replaces occurrences of old with new in a text frame.
    Returns (count_left, num_replacements_made). Preserves run formatting when possible.
    """
    if not old or count_left == 0:
        return count_left, 0
    num_done = 0
    for p in tf.paragraphs:
        if old not in p.text:
            continue
        replaced_in_run = False
        for r in p.runs:
            if old in r.text:
                if count_left == -1:
                    occ = r.text.count(old)
                    r.text = r.text.replace(old, new)
                    num_done += occ
                    replaced_in_run = True
                else:
                    occ = r.text.count(old)
                    to_rep = min(occ, count_left)
                    r.text = r.text.replace(old, new, to_rep)
                    count_left -= to_rep
                    num_done += to_rep
                    replaced_in_run = True
                    if count_left <= 0:
                        break
        if not replaced_in_run and old in p.text:
            if count_left == -1:
                occ = p.text.count(old)
                p.text = p.text.replace(old, new)
                num_done += occ
            else:
                occ = p.text.count(old)
                to_rep = min(occ, count_left)
                p.text = p.text.replace(old, new, to_rep)
                count_left -= to_rep
                num_done += to_rep
        if count_left == 0:
            break
    return count_left, num_done


def _replace_in_slide(slide, old: str, new: str, count_left: int) -> tuple[int, int]:
    total_done = 0
    for shape in _iter_all_shapes(slide.shapes):
        if shape.has_text_frame:
            count_left, done = _replace_in_text_frame(shape.text_frame, old, new, count_left)
            total_done += done
            if count_left == 0:
                break
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        count_left, done = _replace_in_text_frame(cell.text_frame, old, new, count_left)
                        total_done += done
                        if count_left == 0:
                            break
                if count_left == 0:
                    break
        if count_left == 0:
            break
    if count_left != 0 and slide.has_notes_slide:
        count_left, done = _replace_in_text_frame(slide.notes_slide.notes_text_frame, old, new, count_left)
        total_done += done
    return count_left, total_done


def edit_pptx(
    input_path: str,
    operations: list[dict[str, Any]] | str,
    output_path: str | None = None,
) -> dict[str, Any]:
    """Modifica una presentación PowerPoint existente (.pptx) in-place o hacia un nuevo archivo vía python-pptx,
    preservando temas, estilos, gráficos y diapositivas existentes.

    Operaciones soportadas (lista de dicts):
    - {"op": "replace_text", "find"|"old": str, "replace"|"new": str, "slide_index"?: int, "count"?: int}
    - {"op": "set_slide_title", "title"|"new_title": str, "slide_index"?: int, "target_title"?: str}
    - {"op": "add_slide", "layout"?: "title_and_content"|"blank", "title"?: str, "bullets"?: list[str]}
    - {"op": "delete_slide", "slide_index"|"index": int}
    - {"op": "append_bullets", "bullets"|"items": list[str], "slide_index"?: int, "shape_index"?: int}
    - {"op": "set_notes", "notes"|"text": str, "slide_index"?: int}

    Devuelve dict con status ("ok"), path, bytes, fidelity ("high"|"partial"), warnings y operations.
    """
    import json
    from .security import safe_out
    from pptx import Presentation

    input_path = os.path.abspath(os.path.expanduser(str(input_path)))
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Archivo PPTX no encontrado: {input_path}")

    target_out = safe_out(output_path if output_path else input_path)
    prs = Presentation(input_path)

    warnings: list[str] = []
    fidelity = "high"

    if isinstance(operations, str):
        try:
            ops = json.loads(operations)
        except json.JSONDecodeError as exc:
            raise ValueError(f"operations inválido (JSON malformado): {exc}")
    else:
        ops = list(operations or [])

    applied_count = 0

    for i, op in enumerate(ops):
        op_name = str(op.get("op") or op.get("operation") or op.get("action") or "").lower().strip()

        try:
            if op_name in ("replace_text", "replace", "text_replace"):
                old = str(op.get("find") if op.get("find") is not None else op.get("old", ""))
                new = str(op.get("replace") if op.get("replace") is not None else op.get("new", ""))
                if not old:
                    warnings.append(f"Op {i} (replace_text): missing 'find' or 'old' parameter.")
                    continue
                raw_count = op.get("count")
                count_left = int(raw_count) if raw_count is not None else -1

                slide_target = op.get("slide_index") if op.get("slide_index") is not None else op.get("slide_idx", op.get("slide"))
                if slide_target is not None:
                    s_idx = int(slide_target)
                    if s_idx < 0:
                        s_idx = len(prs.slides) + s_idx
                    if 0 <= s_idx < len(prs.slides):
                        target_slides = [prs.slides[s_idx]]
                    else:
                        warnings.append(f"Op {i} (replace_text): slide_index {slide_target} out of range.")
                        fidelity = "partial"
                        continue
                else:
                    target_slides = list(prs.slides)

                total_replacements = 0
                for s in target_slides:
                    count_left, done = _replace_in_slide(s, old, new, count_left)
                    total_replacements += done
                    if count_left == 0:
                        break
                applied_count += 1

            elif op_name in ("set_slide_title", "slide_title", "set_title"):
                new_title = str(op.get("title") if op.get("title") is not None else op.get("new_title", op.get("text", "")))
                target_idx = op.get("slide_index") if op.get("slide_index") is not None else op.get("slide_idx", op.get("slide"))
                target_title = op.get("target_title") or op.get("old_title") or op.get("find_title") or op.get("current_title")

                matched_slide = None
                target_p = None

                if target_title:
                    target_clean = str(target_title).strip().lower()
                    for s in prs.slides:
                        if s.shapes.title and target_clean in s.shapes.title.text.strip().lower():
                            matched_slide = s
                            break
                        for sh in s.shapes:
                            if sh.has_text_frame:
                                for p in sh.text_frame.paragraphs:
                                    if target_clean in p.text.strip().lower():
                                        matched_slide = s
                                        target_p = p
                                        break
                            if matched_slide:
                                break
                        if matched_slide:
                            break
                    if not matched_slide:
                        warnings.append(f"Op {i} (set_slide_title): slide with title matching '{target_title}' not found.")
                        fidelity = "partial"
                        continue
                else:
                    s_idx = int(target_idx if target_idx is not None else 0)
                    if s_idx < 0:
                        s_idx = len(prs.slides) + s_idx
                    if not (0 <= s_idx < len(prs.slides)):
                        warnings.append(f"Op {i} (set_slide_title): slide_index {target_idx} out of range.")
                        fidelity = "partial"
                        continue
                    matched_slide = prs.slides[s_idx]

                if matched_slide.shapes.title and matched_slide.shapes.title.has_text_frame:
                    tf = matched_slide.shapes.title.text_frame
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = new_title
                        for r in tf.paragraphs[0].runs[1:]:
                            r.text = ""
                    else:
                        tf.text = new_title
                elif target_p is not None:
                    if target_p.runs:
                        target_p.runs[0].text = new_title
                        for r in target_p.runs[1:]:
                            r.text = ""
                    else:
                        target_p.text = new_title
                else:
                    candidate_p = None
                    max_font_size = 0.0
                    for sh in matched_slide.shapes:
                        if sh.has_text_frame:
                            for p in sh.text_frame.paragraphs:
                                txt = p.text.strip()
                                if not txt:
                                    continue
                                sz = 0.0
                                if p.font and p.font.size:
                                    sz = float(p.font.size.pt)
                                elif p.runs and p.runs[0].font.size:
                                    sz = float(p.runs[0].font.size.pt)
                                if sz > max_font_size:
                                    max_font_size = sz
                                    candidate_p = p
                                elif candidate_p is None:
                                    candidate_p = p
                    if candidate_p is not None:
                        if candidate_p.runs:
                            candidate_p.runs[0].text = new_title
                            for r in candidate_p.runs[1:]:
                                r.text = ""
                        else:
                            candidate_p.text = new_title
                    else:
                        from pptx.util import Inches, Pt
                        tb = matched_slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.2))
                        p = tb.text_frame.paragraphs[0]
                        p.text = new_title
                        p.font.size = Pt(30)
                        p.font.bold = True
                applied_count += 1

            elif op_name in ("add_slide", "new_slide", "insert_slide"):
                from pptx.util import Inches, Pt
                from pptx.enum.shapes import PP_PLACEHOLDER

                layout_spec = str(op.get("layout", "title_and_content")).lower().strip()
                title = op.get("title")
                bullets = op.get("bullets") or op.get("items") or []
                if isinstance(bullets, str):
                    bullets = [bullets]

                selected_layout = None
                if layout_spec in ("blank", "empty"):
                    for lyt in prs.slide_layouts:
                        if "blank" in lyt.name.lower():
                            selected_layout = lyt
                            break
                    if selected_layout is None:
                        selected_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
                else:
                    for lyt in prs.slide_layouts:
                        if "title and content" in lyt.name.lower() or "title & content" in lyt.name.lower():
                            selected_layout = lyt
                            break
                    if selected_layout is None:
                        for lyt in prs.slide_layouts:
                            types = [ph.type for ph in lyt.placeholders]
                            if any(t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) for t in types):
                                selected_layout = lyt
                                break
                    if selected_layout is None:
                        selected_layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]

                new_s = prs.slides.add_slide(selected_layout)

                if title is not None:
                    if new_s.shapes.title and new_s.shapes.title.has_text_frame:
                        new_s.shapes.title.text = str(title)
                    else:
                        tb = new_s.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.2))
                        p = tb.text_frame.paragraphs[0]
                        p.text = str(title)
                        p.font.size = Pt(30)
                        p.font.bold = True

                if bullets:
                    content_shape = None
                    for ph in new_s.placeholders:
                        if ph == new_s.shapes.title:
                            continue
                        content_shape = ph
                        break

                    if content_shape and content_shape.has_text_frame:
                        tf = content_shape.text_frame
                        for idx_b, b in enumerate(bullets):
                            p = tf.paragraphs[0] if (idx_b == 0 and not tf.paragraphs[0].text) else tf.add_paragraph()
                            p.text = str(b)
                            p.level = 0
                    else:
                        top_pos = Inches(2.2) if title else Inches(1.0)
                        tb = new_s.shapes.add_textbox(Inches(1.0), top_pos, Inches(11.333), Inches(4.5))
                        tf = tb.text_frame
                        tf.word_wrap = True
                        for idx_b, b in enumerate(bullets):
                            p = tf.paragraphs[0] if idx_b == 0 else tf.add_paragraph()
                            p.text = str(b)
                            p.level = 0
                applied_count += 1

            elif op_name in ("delete_slide", "remove_slide", "del_slide"):
                index = op.get("slide_index") if op.get("slide_index") is not None else op.get("index", op.get("slide"))
                if index is None:
                    warnings.append(f"Op {i} (delete_slide): missing 'slide_index' parameter.")
                    fidelity = "partial"
                    continue
                idx = int(index)
                if idx < 0:
                    idx = len(prs.slides) + idx
                if not (0 <= idx < len(prs.slides)):
                    warnings.append(f"Op {i} (delete_slide): slide_index {index} out of range.")
                    fidelity = "partial"
                    continue
                r_id = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(r_id)
                del prs.slides._sldIdLst[idx]
                applied_count += 1

            elif op_name in ("append_bullets", "add_bullets", "append_bullet"):
                bullets = op.get("bullets") or op.get("items") or []
                if isinstance(bullets, str):
                    bullets = [bullets]
                if not bullets:
                    warnings.append(f"Op {i} (append_bullets): missing 'bullets' or 'items' parameter.")
                    continue

                target_idx = op.get("slide_index") if op.get("slide_index") is not None else op.get("slide_idx", op.get("slide", 0))
                s_idx = int(target_idx)
                if s_idx < 0:
                    s_idx = len(prs.slides) + s_idx
                if not (0 <= s_idx < len(prs.slides)):
                    warnings.append(f"Op {i} (append_bullets): slide_index {target_idx} out of range.")
                    fidelity = "partial"
                    continue

                slide = prs.slides[s_idx]
                content_tf = None
                shape_idx = op.get("shape_index")
                if shape_idx is not None and 0 <= int(shape_idx) < len(slide.shapes):
                    target_sh = slide.shapes[int(shape_idx)]
                    if target_sh.has_text_frame:
                        content_tf = target_sh.text_frame

                if content_tf is None:
                    for ph in slide.placeholders:
                        if ph == slide.shapes.title:
                            continue
                        if ph.has_text_frame:
                            content_tf = ph.text_frame
                            break

                if content_tf is None:
                    for sh in slide.shapes:
                        if sh == slide.shapes.title:
                            continue
                        if sh.has_text_frame:
                            for p in sh.text_frame.paragraphs:
                                if "▪" in p.text or "•" in p.text:
                                    content_tf = sh.text_frame
                                    break
                            if content_tf:
                                break

                if content_tf is None:
                    for sh in slide.shapes:
                        if sh != slide.shapes.title and sh.has_text_frame:
                            content_tf = sh.text_frame
                            break

                if content_tf is None:
                    from pptx.util import Inches
                    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
                    content_tf = tb.text_frame
                    content_tf.word_wrap = True

                has_sym_run = False
                ref_sym_font = None
                ref_txt_font = None
                if content_tf.paragraphs:
                    last_p = content_tf.paragraphs[-1]
                    if last_p.runs and len(last_p.runs) >= 2 and ("▪" in last_p.runs[0].text or "•" in last_p.runs[0].text):
                        has_sym_run = True
                        ref_sym_font = last_p.runs[0].font
                        ref_txt_font = last_p.runs[1].font
                    elif last_p.runs:
                        ref_txt_font = last_p.runs[0].font

                for b in bullets:
                    p = content_tf.add_paragraph()
                    if has_sym_run and ref_sym_font:
                        r_sym = p.add_run()
                        r_sym.text = "▪  "
                        if ref_sym_font.size: r_sym.font.size = ref_sym_font.size
                        if ref_sym_font.bold is not None: r_sym.font.bold = ref_sym_font.bold
                        if ref_sym_font.color and ref_sym_font.color.rgb: r_sym.font.color.rgb = ref_sym_font.color.rgb
                        if ref_sym_font.name: r_sym.font.name = ref_sym_font.name

                        r_txt = p.add_run()
                        r_txt.text = str(b)
                        if ref_txt_font:
                            if ref_txt_font.size: r_txt.font.size = ref_txt_font.size
                            if ref_txt_font.color and ref_txt_font.color.rgb: r_txt.font.color.rgb = ref_txt_font.color.rgb
                            if ref_txt_font.name: r_txt.font.name = ref_txt_font.name
                    else:
                        p.text = str(b)
                        p.level = 0
                        if ref_txt_font and p.runs:
                            if ref_txt_font.size: p.runs[0].font.size = ref_txt_font.size
                            if ref_txt_font.color and ref_txt_font.color.rgb: p.runs[0].font.color.rgb = ref_txt_font.color.rgb
                            if ref_txt_font.name: p.runs[0].font.name = ref_txt_font.name
                applied_count += 1

            elif op_name in ("set_notes", "speaker_notes", "notes", "add_notes"):
                target_idx = op.get("slide_index") if op.get("slide_index") is not None else op.get("slide_idx", op.get("slide", 0))
                notes_text = str(op.get("notes") if op.get("notes") is not None else op.get("text", ""))
                s_idx = int(target_idx)
                if s_idx < 0:
                    s_idx = len(prs.slides) + s_idx
                if not (0 <= s_idx < len(prs.slides)):
                    warnings.append(f"Op {i} (set_notes): slide_index {target_idx} out of range.")
                    fidelity = "partial"
                    continue
                slide = prs.slides[s_idx]
                slide.notes_slide.notes_text_frame.text = notes_text
                applied_count += 1

            else:
                warnings.append(f"Op {i}: Operación de edición PPTX no soportada: '{op_name}'")
                fidelity = "partial"

        except Exception as exc:
            warnings.append(f"Op {i} ('{op_name}') falló: {exc}")
            fidelity = "partial"

    prs.save(target_out)
    return {
        "status": "ok",
        "path": os.path.abspath(target_out),
        "bytes": os.path.getsize(target_out),
        "fidelity": fidelity,
        "warnings": warnings,
        "operations": len(ops),
        "operations_applied": applied_count,
        "slides_count": len(prs.slides),
    }

