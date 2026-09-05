"""Extracción estructurada de texto y tablas desde documentos Office (.docx, .pptx, .xlsx, .xlsm)."""
from __future__ import annotations
import os
from pathlib import Path
from typing import Any


def _format_pipe_table(headers: list[Any], rows: list[list[Any]]) -> str:
    """Genera una tabla en formato Markdown con pipes (| ... |)."""
    if not headers and not rows:
        return ""
    if not headers and rows:
        headers = [f"Columna {i+1}" for i in range(len(rows[0]))]

    def clean(val: Any) -> str:
        return str(val if val is not None else "").replace("\n", " ").replace("|", "\\|").strip()

    clean_headers = [clean(h) for h in headers]
    hdr_line = "| " + " | ".join(clean_headers) + " |"
    sep_line = "| " + " | ".join("---" for _ in clean_headers) + " |"
    row_lines = []
    for r in rows:
        clean_row = [clean(c) for c in r]
        if len(clean_row) < len(clean_headers):
            clean_row.extend([""] * (len(clean_headers) - len(clean_row)))
        row_lines.append("| " + " | ".join(clean_row[:len(clean_headers)]) + " |")
    return "\n".join([hdr_line, sep_line] + row_lines)


def read_office(path: str, max_rows: int = 500, format: str = "json") -> dict[str, Any]:
    """Extrae texto estructurado, tablas y diapositivas de archivos Office (.docx, .pptx, .xlsx, .xlsm).

    - path: ruta al archivo de Office.
    - max_rows: límite de filas por hoja al leer planillas Excel (default 500).
    - format: 'json' (default, estructura limpia estructurada) o 'markdown' (texto legible con headings y tablas pipe).

    Devuelve dict con estructura según el formato solicitado.
    """
    path = os.path.abspath(os.path.expanduser(str(path)))
    if not os.path.exists(path):
        raise FileNotFoundError(f"Archivo no encontrado: {path}")

    ext = Path(path).suffix.lower()
    is_markdown = str(format or "").lower().strip() in ("markdown", "md")

    if ext == ".docx":
        from docx import Document
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.table import CT_Tbl
        from docx.text.paragraph import Paragraph
        from docx.table import Table

        doc = Document(path)

        if is_markdown:
            md_blocks = []
            n_paragraphs = 0
            n_tables = 0
            for child in doc.element.body.iterchildren():
                if isinstance(child, CT_P):
                    p = Paragraph(child, doc)
                    txt = p.text.strip()
                    if not txt:
                        continue
                    n_paragraphs += 1
                    st = (p.style.name or "Normal").strip().lower() if p.style else "normal"
                    if st.startswith("heading 1"):
                        md_blocks.append(f"# {txt}")
                    elif st.startswith("heading 2"):
                        md_blocks.append(f"## {txt}")
                    elif st.startswith("heading 3"):
                        md_blocks.append(f"### {txt}")
                    elif st.startswith("heading 4"):
                        md_blocks.append(f"#### {txt}")
                    elif st.startswith("title"):
                        md_blocks.append(f"# {txt}")
                    elif st.startswith("subtitle"):
                        md_blocks.append(f"*{txt}*")
                    elif st.startswith("list") or "bullet" in st:
                        md_blocks.append(f"- {txt}")
                    else:
                        md_blocks.append(txt)
                elif isinstance(child, CT_Tbl):
                    tbl = Table(child, doc)
                    tbl_rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                    if tbl_rows:
                        n_tables += 1
                        headers = tbl_rows[0]
                        data_rows = tbl_rows[1:] if len(tbl_rows) > 1 else []
                        md_blocks.append(_format_pipe_table(headers, data_rows))

            return {
                "status": "ok",
                "path": path,
                "format": "markdown",
                "content": "\n\n".join(md_blocks).strip(),
                "n_paragraphs": n_paragraphs,
                "n_tables": n_tables,
            }

        paragraphs = []
        for p in doc.paragraphs:
            txt = p.text.strip()
            if txt:
                st = p.style.name if p.style else "Normal"
                paragraphs.append({"text": txt, "style": st})

        tables = []
        for tbl in doc.tables:
            tbl_rows = []
            for row in tbl.rows:
                tbl_rows.append([cell.text.strip() for cell in row.cells])
            tables.append(tbl_rows)

        full_text_parts = [p["text"] for p in paragraphs]
        for tbl in tables:
            for r in tbl:
                full_text_parts.append(" | ".join(r))

        return {
            "status": "ok",
            "path": path,
            "format": "docx",
            "n_paragraphs": len(paragraphs),
            "n_tables": len(tables),
            "paragraphs": paragraphs,
            "tables": tables,
            "text": "\n\n".join(full_text_parts).strip(),
        }

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        slides_data = []
        all_text_parts = []
        md_slides = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            slide_tables = []
            slide_md_lines = []
            title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        if not title and (shape == slide.shapes.title or "title" in shape.name.lower()):
                            title = t
                        else:
                            slide_texts.append(t)
                            slide_md_lines.append(f"- {t}")
                        all_text_parts.append(t)
                if shape.has_table:
                    tbl_content = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    slide_tables.append(tbl_content)
                    if tbl_content:
                        headers = tbl_content[0]
                        data_rows = tbl_content[1:] if len(tbl_content) > 1 else []
                        slide_md_lines.append(_format_pipe_table(headers, data_rows))
                    for r in tbl_content:
                        all_text_parts.append(" | ".join(r))

            s_title = title or (slide_texts[0] if slide_texts else f"Diapositiva {idx}")
            slides_data.append({
                "slide": idx,
                "title": s_title,
                "text": slide_texts,
                "tables": slide_tables,
            })

            slide_md = f"## Slide {idx}: {s_title}"
            if slide_md_lines:
                slide_md += "\n\n" + "\n".join(slide_md_lines)
            md_slides.append(slide_md)

        if is_markdown:
            return {
                "status": "ok",
                "path": path,
                "format": "markdown",
                "content": "\n\n".join(md_slides).strip(),
                "n_slides": len(prs.slides),
            }

        return {
            "status": "ok",
            "path": path,
            "format": "pptx",
            "n_slides": len(prs.slides),
            "slides": slides_data,
            "text": "\n\n".join(all_text_parts).strip(),
        }

    elif ext in (".xlsx", ".xlsm"):
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        sheets_data = []
        md_sections = []
        limit_r = int(max_rows) if max_rows else 500

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for r_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                if r_idx > limit_r:
                    break
                if any(c is not None and str(c).strip() != "" for c in row):
                    rows.append([str(c) if c is not None else "" for c in row])

            headers = rows[0] if rows else []
            data_rows = rows[1:] if len(rows) > 1 else []
            sheets_data.append({
                "name": sheet_name,
                "max_row": ws.max_row or len(rows),
                "max_column": ws.max_column or (len(headers) if headers else 0),
                "headers": headers,
                "rows": data_rows,
            })

            sec_md = f"## {sheet_name}\n\n"
            if rows:
                sec_md += _format_pipe_table(headers, data_rows)
            else:
                sec_md += "*(Hoja vacía)*"
            md_sections.append(sec_md)

        wb.close()

        if is_markdown:
            return {
                "status": "ok",
                "path": path,
                "format": "markdown",
                "content": "\n\n".join(md_sections).strip(),
                "n_sheets": len(sheets_data),
            }

        return {
            "status": "ok",
            "path": path,
            "format": "xlsx",
            "n_sheets": len(sheets_data),
            "sheets": sheets_data,
        }

    else:
        raise ValueError(f"Formato no soportado para lectura estructurada: '{ext}'. Formatos válidos: .docx, .pptx, .xlsx, .xlsm")
